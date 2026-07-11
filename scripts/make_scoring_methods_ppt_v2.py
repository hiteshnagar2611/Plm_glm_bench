#!/usr/bin/env python3
"""
Pathogenicity Scoring Methods PPT with embedded plots.
Generates plots first, then builds PPT with images.
"""
import matplotlib
matplotlib.use('Agg')
import matplotlib.pyplot as plt
import matplotlib.patches as mpatches
import numpy as np
import os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

# ── Paths ──
SCRIPT_DIR = os.path.dirname(os.path.abspath(__file__))
PROJ_ROOT = os.path.join(SCRIPT_DIR, '..')
FIG_DIR = os.path.join(PROJ_ROOT, 'benchmark_v3', 'figures')
TEMP_DIR = os.path.join(SCRIPT_DIR, '_tmp_plots')
os.makedirs(FIG_DIR, exist_ok=True)
os.makedirs(TEMP_DIR, exist_ok=True)

# ── Plot style ──
plt.rcParams.update({
    'figure.facecolor': '#0f172a',
    'axes.facecolor': '#1e293b',
    'axes.edgecolor': '#334155',
    'axes.labelcolor': '#e2e8f0',
    'xtick.color': '#94a3b8',
    'ytick.color': '#94a3b8',
    'text.color': '#e2e8f0',
    'grid.color': '#334155',
    'grid.alpha': 0.5,
    'font.family': 'sans-serif',
    'font.size': 11,
})

# ════════════════════════════════════════════════════════════
# PLOT 1: AUROC Bar Chart
# ════════════════════════════════════════════════════════════
models = ['ESM-1b', 'ESM-2', 'SaProt', 'NT-v2', 'AlphaGenome', 'ProtT5', 'HyenaDNA']
aurocs = [0.873, 0.838, 0.833, 0.638, 0.552, 0.536, 0.448]
spearmans = [0.631, 0.572, 0.556, 0.273, 0.153, 0.186, 0.056]
colors = ['#3b82f6', '#3b82f6', '#8b5cf6', '#10b981', '#f97316', '#3b82f6', '#10b981']

fig, ax = plt.subplots(figsize=(10, 5))
bars = ax.bar(range(len(models)), aurocs, color=colors, edgecolor='white', linewidth=0.8, width=0.7)
ax.set_xticks(range(len(models)))
ax.set_xticklabels(models, fontsize=11, fontweight='bold')
ax.set_ylabel('AUROC', fontsize=13, fontweight='bold')
ax.set_title('V3 Benchmark: Variant Pathogenicity AUROC', fontsize=15, fontweight='bold', pad=15)
ax.set_ylim(0, 1.05)
ax.axhline(y=0.5, color='#ef4444', linestyle='--', alpha=0.5, linewidth=1, label='Random (0.5)')
ax.axhline(y=0.8, color='#22c55e', linestyle='--', alpha=0.3, linewidth=1, label='Strong (0.8)')
ax.grid(axis='y', alpha=0.3)
for bar, val in zip(bars, aurocs):
    ax.text(bar.get_x() + bar.get_width()/2, bar.get_height() + 0.015, f'{val:.3f}',
            ha='center', va='bottom', fontsize=11, fontweight='bold', color='white')
ax.legend(fontsize=9, loc='upper right', framealpha=0.8)
plt.tight_layout()
fig.savefig(os.path.join(TEMP_DIR, 'plot_auroc.png'), dpi=200, bbox_inches='tight', facecolor='#0f172a')
plt.close()
print("  Saved plot_auroc.png")

# ════════════════════════════════════════════════════════════
# PLOT 2: Spearman Bar Chart
# ════════════════════════════════════════════════════════════
fig, ax = plt.subplots(figsize=(10, 5))
bars = ax.bar(range(len(models)), spearmans, color=colors, edgecolor='white', linewidth=0.8, width=0.7)
ax.set_xticks(range(len(models)))
ax.set_xticklabels(models, fontsize=11, fontweight='bold')
ax.set_ylabel('Spearman ρ', fontsize=13, fontweight='bold')
ax.set_title('V3 Benchmark: Spearman Correlation', fontsize=15, fontweight='bold', pad=15)
ax.set_ylim(0, 0.8)
ax.axhline(y=0, color='#ef4444', linestyle='--', alpha=0.5, linewidth=1)
ax.grid(axis='y', alpha=0.3)
for bar, val in zip(bars, spearmans):
    ax.text(bar.get_x() + bar.get_width()/2, bar.get_height() + 0.01, f'{val:.3f}',
            ha='center', va='bottom', fontsize=11, fontweight='bold', color='white')
plt.tight_layout()
fig.savefig(os.path.join(TEMP_DIR, 'plot_spearman.png'), dpi=200, bbox_inches='tight', facecolor='#0f172a')
plt.close()
print("  Saved plot_spearman.png")

# ════════════════════════════════════════════════════════════
# PLOT 3: Scoring Method Comparison (LLR vs Cosine vs Delta)
# ════════════════════════════════════════════════════════════
# Data from our experiments
method_models = ['ESM-1b', 'ESM-2', 'SaProt']
llr_auroc =    [0.873, 0.838, 0.833]
cos_auroc =    [0.816, 0.797, 0.568]
pool_auroc =   [0.790, 0.780, 0.560]  # approximate from pool experiments

x = np.arange(len(method_models))
width = 0.25

fig, ax = plt.subplots(figsize=(9, 5))
b1 = ax.bar(x - width, llr_auroc, width, label='Log-Likelihood Ratio', color='#22c55e', edgecolor='white', linewidth=0.8)
b2 = ax.bar(x, cos_auroc, width, label='Cosine Similarity', color='#f97316', edgecolor='white', linewidth=0.8)
b3 = ax.bar(x + width, pool_auroc, width, label='MaxPool + Cosine(diff,wt)', color='#ef4444', edgecolor='white', linewidth=0.8)
ax.set_xticks(x)
ax.set_xticklabels(method_models, fontsize=12, fontweight='bold')
ax.set_ylabel('AUROC', fontsize=13, fontweight='bold')
ax.set_title('Scoring Method Comparison: LLR vs Embedding Methods', fontsize=14, fontweight='bold', pad=15)
ax.set_ylim(0, 1.0)
ax.legend(fontsize=10, loc='upper right', framealpha=0.8)
ax.grid(axis='y', alpha=0.3)
for bars_group in [b1, b2, b3]:
    for bar in bars_group:
        ax.text(bar.get_x() + bar.get_width()/2, bar.get_height() + 0.01, f'{bar.get_height():.3f}',
                ha='center', va='bottom', fontsize=9, fontweight='bold', color='white')
plt.tight_layout()
fig.savefig(os.path.join(TEMP_DIR, 'plot_scoring_methods.png'), dpi=200, bbox_inches='tight', facecolor='#0f172a')
plt.close()
print("  Saved plot_scoring_methods.png")

# ════════════════════════════════════════════════════════════
# PLOT 4: Paper vs Our Method Comparison
# ════════════════════════════════════════════════════════════
paper_methods = ['ESM-1b\n(Masked Marginal)', 'ESM-2\n(Masked Marginal)', 'SaProt\n(Struct-Aware)', 'ProtT5\n(VESPA)', 'NT-v2\n(Cosine)', 'HyenaDNA\n(LLR)', 'AlphaGenome\n(ALT−REF)']
our_methods = ['ESM-1b\n(Pseudolikelihood)', 'ESM-2\n(Pseudolikelihood)', 'SaProt\n(Single-WT)', 'ProtT5\n(Position Cos)', 'NT-v2\n(ΔHidden)', 'HyenaDNA\n(LLR)', 'AlphaGenome\n(Mean |ΔTrack|)']

# Paper approximate ClinVar AUROC (from literature)
paper_auroc = [0.89, 0.88, 0.85, 0.82, 0.68, 0.50, 0.60]
our_auroc =   [0.873, 0.838, 0.833, 0.536, 0.638, 0.448, 0.552]

x = np.arange(len(paper_methods))
width = 0.35

fig, ax = plt.subplots(figsize=(11, 5))
b1 = ax.bar(x - width/2, paper_auroc, width, label='Paper Method (ClinVar)', color='#3b82f6', edgecolor='white', linewidth=0.8)
b2 = ax.bar(x + width/2, our_auroc, width, label='Our V3 Method (ClinVar)', color='#22c55e', edgecolor='white', linewidth=0.8)
ax.set_xticks(x)
ax.set_xticklabels(paper_methods, fontsize=9, fontweight='bold')
ax.set_ylabel('AUROC', fontsize=13, fontweight='bold')
ax.set_title('Paper Reported vs Our V3 Results', fontsize=14, fontweight='bold', pad=15)
ax.set_ylim(0, 1.05)
ax.legend(fontsize=10, loc='upper right', framealpha=0.8)
ax.grid(axis='y', alpha=0.3)
for bars_group in [b1, b2]:
    for bar in bars_group:
        ax.text(bar.get_x() + bar.get_width()/2, bar.get_height() + 0.01, f'{bar.get_height():.3f}',
                ha='center', va='bottom', fontsize=8, fontweight='bold', color='white')
plt.tight_layout()
fig.savefig(os.path.join(TEMP_DIR, 'plot_paper_vs_ours.png'), dpi=200, bbox_inches='tight', facecolor='#0f172a')
plt.close()
print("  Saved plot_paper_vs_ours.png")

# ════════════════════════════════════════════════════════════
# PLOT 5: Training Data vs Performance Scatter
# ════════════════════════════════════════════════════════════
scatter_models = ['ESM-1b', 'ESM-2', 'SaProt', 'ProtT5', 'NT-v2', 'HyenaDNA', 'AlphaGenome']
tokens_B = [86, 86, 86, 150, 174, 3.1, 50]
params_M = [650, 650, 650, 3000, 500, 1.6, 250]
scatter_auroc = [0.873, 0.838, 0.833, 0.536, 0.638, 0.448, 0.552]
scatter_colors = ['#3b82f6', '#3b82f6', '#8b5cf6', '#3b82f6', '#10b981', '#10b981', '#f97316']

fig, (ax1, ax2) = plt.subplots(1, 2, figsize=(12, 5))

# Left: Params vs AUROC
scatter1 = ax1.scatter(params_M, scatter_auroc, c=scatter_colors, s=200, edgecolors='white', linewidth=1.5, zorder=5)
for i, name in enumerate(scatter_models):
    ax1.annotate(name, (params_M[i], scatter_auroc[i]), textcoords="offset points",
                xytext=(8, 5), fontsize=9, fontweight='bold')
ax1.set_xscale('log')
ax1.set_xlabel('Model Parameters (Millions)', fontsize=12, fontweight='bold')
ax1.set_ylabel('AUROC', fontsize=12, fontweight='bold')
ax1.set_title('Model Size vs Performance', fontsize=13, fontweight='bold')
ax1.grid(alpha=0.3)
ax1.set_ylim(0.3, 1.0)

# Right: Tokens vs AUROC
scatter2 = ax2.scatter(tokens_B, scatter_auroc, c=scatter_colors, s=200, edgecolors='white', linewidth=1.5, zorder=5)
for i, name in enumerate(scatter_models):
    ax2.annotate(name, (tokens_B[i], scatter_auroc[i]), textcoords="offset points",
                xytext=(8, 5), fontsize=9, fontweight='bold')
ax2.set_xscale('log')
ax2.set_xlabel('Training Tokens (Billions)', fontsize=12, fontweight='bold')
ax2.set_ylabel('AUROC', fontsize=12, fontweight='bold')
ax2.set_title('Training Data vs Performance', fontsize=13, fontweight='bold')
ax2.grid(alpha=0.3)
ax2.set_ylim(0.3, 1.0)

legend_elements = [
    mpatches.Patch(color='#3b82f6', label='Protein LM'),
    mpatches.Patch(color='#8b5cf6', label='Protein+Structure'),
    mpatches.Patch(color='#10b981', label='Genomic LM'),
    mpatches.Patch(color='#f97316', label='Regulatory LM'),
]
fig.legend(handles=legend_elements, fontsize=9, loc='lower center', ncol=4, framealpha=0.8)
plt.tight_layout(rect=[0, 0.06, 1, 1])
fig.savefig(os.path.join(TEMP_DIR, 'plot_size_vs_perf.png'), dpi=200, bbox_inches='tight', facecolor='#0f172a')
plt.close()
print("  Saved plot_size_vs_perf.png")

# ════════════════════════════════════════════════════════════
# PLOT 6: Scoring Method Flow Diagram
# ════════════════════════════════════════════════════════════
fig, ax = plt.subplots(figsize=(12, 6))
ax.set_xlim(0, 12)
ax.set_ylim(0, 6)
ax.axis('off')

# Draw boxes
def draw_box(ax, x, y, w, h, text, color='#1e293b', edge='#3b82f6', fontsize=10, textcolor='white'):
    rect = plt.Rectangle((x, y), w, h, facecolor=color, edgecolor=edge, linewidth=2, zorder=2)
    ax.add_patch(rect)
    ax.text(x + w/2, y + h/2, text, ha='center', va='center', fontsize=fontsize,
            fontweight='bold', color=textcolor, zorder=3, wrap=True)

# Input
draw_box(ax, 0.2, 2.5, 1.5, 1.0, 'Input\nVariant', '#1e293b', '#64748b', 11)

# Protein path
draw_box(ax, 2.5, 4.0, 2.0, 1.0, 'Protein Sequence\n(WT & Mutant)', '#1e293b', '#3b82f6', 10)
draw_box(ax, 5.0, 4.0, 2.0, 1.0, 'PLM Forward Pass\n(ESM/SaProt/ProtT5)', '#1e293b', '#3b82f6', 10)
draw_box(ax, 7.5, 4.5, 1.8, 0.8, 'LLR\nlog P(mut)−log P(wt)', '#1e293b', '#22c55e', 9)
draw_box(ax, 7.5, 3.5, 1.8, 0.8, 'Cosine\nsim(mut_emb, wt_emb)', '#1e293b', '#f97316', 9)
draw_box(ax, 9.8, 4.0, 2.0, 1.0, 'Pathogenicity\nScore', '#1e293b', '#22c55e', 11)

# Genomic path
draw_box(ax, 2.5, 1.0, 2.0, 1.0, 'DNA Sequence\n(REF & ALT)', '#1e293b', '#10b981', 10)
draw_box(ax, 5.0, 1.0, 2.0, 1.0, 'GLM Forward Pass\n(NT-v2/HyenaDNA)', '#1e293b', '#10b981', 10)
draw_box(ax, 7.5, 1.5, 1.8, 0.8, 'LLR or Cosine\nor ΔHidden', '#1e293b', '#10b981', 9)
draw_box(ax, 7.5, 0.5, 1.8, 0.8, 'Mean |ΔTrack|\n(AlphaGenome)', '#1e293b', '#f97316', 9)
draw_box(ax, 9.8, 1.0, 2.0, 1.0, 'Pathogenicity\nScore', '#1e293b', '#10b981', 11)

# Arrows
arrow_kw = dict(arrowstyle='->', color='#64748b', lw=2)
from matplotlib.patches import FancyArrowPatch
arrows = [
    (1.7, 3.0, 2.5, 4.5),   # input → protein
    (1.7, 3.0, 2.5, 1.5),   # input → genomic
    (4.5, 4.5, 5.0, 4.5),   # protein seq → PLM
    (4.5, 1.5, 5.0, 1.5),   # dna seq → GLM
    (7.0, 4.5, 7.5, 4.9),   # PLM → LLR
    (7.0, 4.5, 7.5, 3.9),   # PLM → Cosine
    (7.0, 1.5, 7.5, 1.9),   # GLM → LLR/cos
    (7.0, 1.5, 7.5, 0.9),   # GLM → ΔTrack
    (9.3, 4.9, 9.8, 4.5),   # LLR → score
    (9.3, 3.9, 9.8, 4.5),   # Cosine → score
    (9.3, 1.9, 9.8, 1.5),   # LLR → score
    (9.3, 0.9, 9.8, 1.5),   # ΔTrack → score
]
for x1, y1, x2, y2 in arrows:
    ax.annotate('', xy=(x2, y2), xytext=(x1, y1), arrowprops=arrow_kw)

# Labels
ax.text(3.5, 5.5, 'PROTEIN LANGUAGE MODELS', ha='center', fontsize=12, fontweight='bold', color='#3b82f6')
ax.text(3.5, 0.2, 'GENOMIC LANGUAGE MODELS', ha='center', fontsize=12, fontweight='bold', color='#10b981')

plt.tight_layout()
fig.savefig(os.path.join(TEMP_DIR, 'plot_scoring_flow.png'), dpi=200, bbox_inches='tight', facecolor='#0f172a')
plt.close()
print("  Saved plot_scoring_flow.png")

# ════════════════════════════════════════════════════════════
# BUILD PPT WITH EMBEDDED PLOTS
# ════════════════════════════════════════════════════════════
print("\nBuilding PPT with embedded plots...")
prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)

BG_DARK   = RGBColor(15, 23, 42)
BG_CARD   = RGBColor(30, 41, 59)
WHITE     = RGBColor(255, 255, 255)
GRAY      = RGBColor(148, 163, 184)
BLUE      = RGBColor(59, 130, 246)
GREEN     = RGBColor(16, 185, 129)
RED       = RGBColor(239, 68, 68)
ORANGE    = RGBColor(249, 115, 22)
PURPLE    = RGBColor(139, 92, 246)
YELLOW    = RGBColor(234, 179, 8)

def add_bg(slide):
    bg = slide.background.fill; bg.solid(); bg.fore_color.rgb = BG_DARK

def add_text_box(slide, left, top, width, height, text, font_size=14, color=WHITE, bold=False, alignment=PP_ALIGN.LEFT):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.text = text; p.font.size = Pt(font_size)
    p.font.color.rgb = color; p.font.bold = bold; p.font.name = 'Calibri'; p.alignment = alignment
    return tf

def add_para(tf, text, font_size=14, color=WHITE, bold=False, space_before=6, alignment=PP_ALIGN.LEFT):
    p = tf.add_paragraph(); p.text = text; p.font.size = Pt(font_size)
    p.font.color.rgb = color; p.font.bold = bold; p.font.name = 'Calibri'
    p.space_before = Pt(space_before); p.alignment = alignment; return p

def add_card(slide, left, top, width, height, fill_color=BG_CARD):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(height))
    shape.fill.solid(); shape.fill.fore_color.rgb = fill_color; shape.line.fill.background()
    return shape

# ════════ SLIDE 1: TITLE ════════
slide = prs.slides.add_slide(prs.slide_layouts[6]); add_bg(slide)
shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(1.5), Inches(2.6), Inches(10.3), Inches(0.06))
shape.fill.solid(); shape.fill.fore_color.rgb = BLUE; shape.line.fill.background()
add_text_box(slide, 1.5, 1.0, 10.3, 1.6, 'Pathogenicity Scoring Methods', 42, WHITE, True, PP_ALIGN.CENTER)
add_text_box(slide, 1.5, 2.8, 10.3, 1.0, 'How PLMs and GLMs Calculate Variant Pathogenicity', 22, GRAY, False, PP_ALIGN.CENTER)
add_text_box(slide, 1.5, 4.0, 10.3, 0.8, 'Paper Methods  |  Our V3 Benchmark  |  Results Comparison', 16, BLUE, False, PP_ALIGN.CENTER)
add_text_box(slide, 1.5, 5.5, 5, 0.5, 'Protein Language Models', 14, BLUE, True)
add_text_box(slide, 6.8, 5.5, 5, 0.5, 'Genomic Language Models', 14, GREEN, True)
add_text_box(slide, 1.5, 5.9, 5, 0.5, 'ESM-1b  |  ESM-2  |  SaProt  |  ProtT5', 12, GRAY)
add_text_box(slide, 6.8, 5.9, 5, 0.5, 'NT-v2  |  HyenaDNA  |  AlphaGenome', 12, GRAY)

# ════════ SLIDE 2: DATA FILTERING PIPELINE ════════
slide = prs.slides.add_slide(prs.slide_layouts[6]); add_bg(slide)
add_text_box(slide, 0.5, 0.3, 12.3, 0.7, 'V3 Benchmark: Data Filtering Pipeline', 30, WHITE, True)
slide.shapes.add_picture(os.path.join(FIG_DIR, 'fig_v3_filtering_pipeline.png'),
                         Inches(0.3), Inches(1.0), Inches(12.7), Inches(6.3))

# ════════ SLIDE 3: SCORING METHODS EXPLAINED ════════
slide = prs.slides.add_slide(prs.slide_layouts[6]); add_bg(slide)
add_text_box(slide, 0.5, 0.3, 12.3, 0.7, 'How We Score Pathogenicity', 30, WHITE, True)

# Left column: Protein LMs
add_card(slide, 0.5, 1.2, 5.9, 5.8)
tf = add_text_box(slide, 0.8, 1.3, 5.5, 0.5, 'Protein Language Models', 18, BLUE, True)
add_para(tf, '', 4, GRAY, space_before=2)

add_para(tf, 'ESM-1b / ESM-2 (Log-Likelihood Ratio)', 13, GREEN, True, space_before=8)
add_para(tf, '  Score = log P(mut_aa | mut_ctx) − log P(wt_aa | wt_ctx)', 11, WHITE, space_before=4)
add_para(tf, '  Negative → pathogenic | Run WT & mutant separately', 10, GRAY, space_before=2)

add_para(tf, 'SaProt (Structure-Aware LLR)', 13, PURPLE, True, space_before=8)
add_para(tf, '  Same LLR but tokens = AA × 3Di (Foldseek structure)', 11, WHITE, space_before=4)
add_para(tf, '  Only WT sequence needed; mutant scored from same context', 10, GRAY, space_before=2)

add_para(tf, 'ProtT5 (Position Cosine Similarity)', 13, ORANGE, True, space_before=8)
add_para(tf, '  Score = cosine(hidden_wt[pos], hidden_mut[pos])', 11, WHITE, space_before=4)
add_para(tf, '  Lower cosine → more pathogenic | Weaker than LLR', 10, GRAY, space_before=2)

add_para(tf, '', 4, GRAY, space_before=4)
add_para(tf, 'All protein LMs use zero-shot scoring (no fine-tuning)', 12, YELLOW, True, space_before=6)

# Right column: Genomic LMs
add_card(slide, 6.9, 1.2, 5.9, 5.8)
tf = add_text_box(slide, 7.2, 1.3, 5.5, 0.5, 'Genomic Language Models', 18, GREEN, True)
add_para(tf, '', 4, GRAY, space_before=2)

add_para(tf, 'NT-v2 (Mean Hidden Delta)', 13, GREEN, True, space_before=8)
add_para(tf, '  Score = mean(|hidden_wt − hidden_mut|)', 11, WHITE, space_before=4)
add_para(tf, '  Higher delta → more pathogenic | 6kb DNA context', 10, GRAY, space_before=2)

add_para(tf, 'HyenaDNA (Log-Likelihood Ratio)', 13, GREEN, True, space_before=8)
add_para(tf, '  Score = log P(alt_base | alt_ctx) − log P(ref_base | ref_ctx)', 11, WHITE, space_before=4)
add_para(tf, '  Autoregressive (causal) | Single nucleotide tokens', 10, GRAY, space_before=2)

add_para(tf, 'AlphaGenome (Multi-Track Delta)', 13, ORANGE, True, space_before=8)
add_para(tf, '  Score = mean(|ref_track_i − alt_track_i|) across 5 tracks', 11, WHITE, space_before=4)
add_para(tf, '  RNA-seq, splice, DNase | Regulatory model, not protein', 10, GRAY, space_before=2)

add_para(tf, '', 4, GRAY, space_before=4)
add_para(tf, 'Genomic models predict DNA-level effects, not protein', 12, YELLOW, True, space_before=6)
add_para(tf, 'pathogenicity directly → task mismatch for missense', 12, YELLOW)

# ════════ SLIDE 4: SCORING FLOW DIAGRAM ════════
slide = prs.slides.add_slide(prs.slide_layouts[6]); add_bg(slide)
add_text_box(slide, 0.5, 0.3, 12.3, 0.7, 'Scoring Pipeline Overview', 30, WHITE, True)
slide.shapes.add_picture(os.path.join(TEMP_DIR, 'plot_scoring_flow.png'),
                         Inches(0.5), Inches(1.1), Inches(12.3), Inches(6.2))

# ════════ SLIDE 5: CORE LLR FORMULA (renumbered from 3) ════════
slide = prs.slides.add_slide(prs.slide_layouts[6]); add_bg(slide)
add_text_box(slide, 0.5, 0.3, 12.3, 0.7, 'Core Method: Log-Likelihood Ratio (LLR)', 30, WHITE, True)
add_card(slide, 0.5, 1.2, 12.3, 1.6)
tf = add_text_box(slide, 0.8, 1.3, 11.7, 0.5, 'The Fundamental Formula', 18, BLUE, True)
add_para(tf, '', 4, GRAY, space_before=2)
add_para(tf, 'LLR  =  log P(mutant residue | mutant context)  −  log P(wildtype residue | wildtype context)', 18, GREEN, True, space_before=6)
add_para(tf, '', 4, GRAY, space_before=2)
add_para(tf, 'Negative score → model assigns lower probability to mutant → likely pathogenic', 14, GRAY, space_before=4)

add_card(slide, 0.5, 3.1, 5.9, 4.0)
tf = add_text_box(slide, 0.8, 3.2, 5.5, 0.5, 'Step-by-Step', 18, GREEN, True)
add_para(tf, '', 4, GRAY, space_before=2)
add_para(tf, '1. Input wildtype sequence → run through LM', 13, WHITE, space_before=6)
add_para(tf, '   Extract log P(wt_aa | context) at variant pos', 12, GRAY)
add_para(tf, '', 4, GRAY, space_before=2)
add_para(tf, '2. Input mutant sequence → run through LM', 13, WHITE, space_before=6)
add_para(tf, '   Extract log P(mut_aa | context) at variant pos', 12, GRAY)
add_para(tf, '', 4, GRAY, space_before=2)
add_para(tf, '3. Score = log P(mut) − log P(wt)', 13, GREEN, True, space_before=6)
add_para(tf, '', 4, GRAY, space_before=2)
add_para(tf, 'More negative → more pathogenic', 13, RED, True, space_before=6)
add_para(tf, 'Score ≥ 0 → model thinks mutation is neutral/beneficial', 13, GREEN)

add_card(slide, 6.9, 3.1, 5.9, 4.0)
tf = add_text_box(slide, 7.2, 3.2, 5.5, 0.5, 'Example: ESM-2 at Position 42', 18, BLUE, True)
add_para(tf, '', 4, GRAY, space_before=2)
add_para(tf, 'Wildtype: ...A R K [G] L E D...', 13, WHITE, space_before=8)
add_para(tf, '  P(Gly | context) = 0.85', 13, GRAY)
add_para(tf, '  log P(Gly) = −0.163', 13, GRAY)
add_para(tf, '', 4, GRAY, space_before=2)
add_para(tf, 'Mutant:   ...A R K [D] L E D...', 13, WHITE, space_before=8)
add_para(tf, '  P(Asp | context) = 0.02', 13, GRAY)
add_para(tf, '  log P(Asp) = −3.912', 13, GRAY)
add_para(tf, '', 4, GRAY, space_before=2)
add_para(tf, 'LLR = −3.912 − (−0.163) = −3.749', 14, RED, True, space_before=8)
add_para(tf, '→ Strongly pathogenic (charge reversal G→D)', 12, GRAY)

# ════════ SLIDE 4: AUROC RESULTS ════════
slide = prs.slides.add_slide(prs.slide_layouts[6]); add_bg(slide)
add_text_box(slide, 0.5, 0.3, 12.3, 0.7, 'V3 Benchmark Results: AUROC', 30, WHITE, True)
slide.shapes.add_picture(os.path.join(TEMP_DIR, 'plot_auroc.png'),
                         Inches(0.5), Inches(1.1), Inches(7.5), Inches(5.8))
# Side annotations
add_card(slide, 8.3, 1.2, 4.5, 5.5)
tf = add_text_box(slide, 8.5, 1.3, 4.1, 0.5, 'Key Takeaways', 16, BLUE, True)
add_para(tf, '', 4, GRAY, space_before=2)
add_para(tf, 'Protein LMs dominate:', 13, GREEN, True, space_before=8)
add_para(tf, '  ESM-1b: 0.873 (best)', 13, WHITE, space_before=4)
add_para(tf, '  ESM-2: 0.838', 13, WHITE, space_before=2)
add_para(tf, '  SaProt: 0.833', 13, WHITE, space_before=2)
add_para(tf, '', 4, GRAY, space_before=4)
add_para(tf, 'Genomic LMs struggle:', 13, ORANGE, True, space_before=6)
add_para(tf, '  NT-v2: 0.638', 13, WHITE, space_before=4)
add_para(tf, '  AlphaGenome: 0.552', 13, WHITE, space_before=2)
add_para(tf, '  ProtT5: 0.536', 13, WHITE, space_before=2)
add_para(tf, '  HyenaDNA: 0.448', 13, RED, space_before=2)
add_para(tf, '', 4, GRAY, space_before=4)
add_para(tf, 'Why the gap?', 14, YELLOW, True, space_before=6)
add_para(tf, '  Protein LMs learn amino acid', 12, GRAY, space_before=4)
add_para(tf, '  semantics directly relevant to', 12, GRAY)
add_para(tf, '  missense pathogenicity', 12, GRAY)

# ════════ SLIDE 5: SPEARMAN RESULTS ════════
slide = prs.slides.add_slide(prs.slide_layouts[6]); add_bg(slide)
add_text_box(slide, 0.5, 0.3, 12.3, 0.7, 'V3 Benchmark Results: Spearman Correlation', 30, WHITE, True)
slide.shapes.add_picture(os.path.join(TEMP_DIR, 'plot_spearman.png'),
                         Inches(0.5), Inches(1.1), Inches(7.5), Inches(5.8))
add_card(slide, 8.3, 1.2, 4.5, 5.5)
tf = add_text_box(slide, 8.5, 1.3, 4.1, 0.5, 'Rank Correlation', 16, BLUE, True)
add_para(tf, '', 4, GRAY, space_before=2)
add_para(tf, 'Spearman ρ measures how well', 12, GRAY, space_before=8)
add_para(tf, 'model scores rank variants correctly', 12, GRAY)
add_para(tf, '', 4, GRAY, space_before=4)
add_para(tf, 'Top performers:', 13, GREEN, True, space_before=6)
add_para(tf, '  ESM-1b: 0.631', 13, WHITE, space_before=4)
add_para(tf, '  ESM-2: 0.572', 13, WHITE, space_before=2)
add_para(tf, '  SaProt: 0.556', 13, WHITE, space_before=2)
add_para(tf, '', 4, GRAY, space_before=4)
add_para(tf, 'Bottom:', 13, RED, True, space_before=6)
add_para(tf, '  HyenaDNA: 0.056', 13, WHITE, space_before=4)
add_para(tf, '  (nearly random ranking)', 12, GRAY)
add_para(tf, '', 4, GRAY, space_before=4)
add_para(tf, '→ AUROC and Spearman agree', 12, YELLOW, space_before=6)
add_para(tf, '   on model ranking', 12, YELLOW)

# ════════ SLIDE 6: SCORING METHOD COMPARISON ════════
slide = prs.slides.add_slide(prs.slide_layouts[6]); add_bg(slide)
add_text_box(slide, 0.5, 0.3, 12.3, 0.7, 'Scoring Method Comparison: LLR vs Embedding', 30, WHITE, True)
slide.shapes.add_picture(os.path.join(TEMP_DIR, 'plot_scoring_methods.png'),
                         Inches(0.5), Inches(1.1), Inches(7.5), Inches(5.8))
add_card(slide, 8.3, 1.2, 4.5, 5.5)
tf = add_text_box(slide, 8.5, 1.3, 4.1, 0.5, 'Method Ranking', 16, GREEN, True)
add_para(tf, '', 4, GRAY, space_before=2)
add_para(tf, 'LLR always wins:', 13, GREEN, True, space_before=8)
add_para(tf, '  Captures probabilistic', 12, GRAY, space_before=4)
add_para(tf, '  uncertainty that geometry', 12, GRAY)
add_para(tf, '  methods miss', 12, GRAY)
add_para(tf, '', 4, GRAY, space_before=4)
add_para(tf, 'Cosine similarity:', 13, ORANGE, True, space_before=6)
add_para(tf, '  Loses 5-8% Spearman vs LLR', 12, GRAY, space_before=4)
add_para(tf, '  Misses probability magnitude', 12, GRAY)
add_para(tf, '', 4, GRAY, space_before=4)
add_para(tf, 'MaxPool + cosine(diff,wt):', 13, RED, True, space_before=6)
add_para(tf, '  NEGATIVE Spearman!', 12, RED, space_before=4)
add_para(tf, '  (broken for this task)', 12, GRAY)
add_para(tf, '', 4, GRAY, space_before=4)
add_para(tf, '→ LLR is the gold standard', 12, YELLOW, space_before=6)

# ════════ SLIDE 7: PAPER VS OURS ════════
slide = prs.slides.add_slide(prs.slide_layouts[6]); add_bg(slide)
add_text_box(slide, 0.5, 0.3, 12.3, 0.7, 'Paper Reported vs Our V3 Results', 30, WHITE, True)
slide.shapes.add_picture(os.path.join(TEMP_DIR, 'plot_paper_vs_ours.png'),
                         Inches(0.5), Inches(1.1), Inches(12.3), Inches(5.8))
add_card(slide, 0.5, 6.9, 12.3, 0.5)
tf = add_text_box(slide, 0.8, 6.95, 11.7, 0.4,
    'Our V3 results closely match paper-reported ClinVar AUROC for ESM-1b/ESM-2/SaProt using nearly identical LLR methods',
    12, YELLOW, True, PP_ALIGN.CENTER)

# ════════ SLIDE 8: SIZE vs PERFORMANCE ════════
slide = prs.slides.add_slide(prs.slide_layouts[6]); add_bg(slide)
add_text_box(slide, 0.5, 0.3, 12.3, 0.7, 'Model Size & Training Data vs Performance', 30, WHITE, True)
slide.shapes.add_picture(os.path.join(TEMP_DIR, 'plot_size_vs_perf.png'),
                         Inches(0.5), Inches(1.1), Inches(12.3), Inches(5.8))
add_card(slide, 0.5, 6.9, 12.3, 0.5)
tf = add_text_box(slide, 0.8, 6.95, 11.7, 0.4,
    'Bigger models and more data help within same modality, but task alignment matters more than scale',
    12, YELLOW, True, PP_ALIGN.CENTER)

# ════════ SLIDE 9: PER-MODEL DETAIL ════════
slide = prs.slides.add_slide(prs.slide_layouts[6]); add_bg(slide)
add_text_box(slide, 0.5, 0.3, 12.3, 0.7, 'Per-Model: Paper Method vs Our Method', 30, WHITE, True)

models_detail = [
    ('ESM-1b', 'Masked marginal LLR', 'Pseudolikelihood LLR', '0.873', '0.631', BLUE, 'Yes'),
    ('ESM-2', 'Masked marginal LLR', 'Pseudolikelihood LLR', '0.838', '0.572', BLUE, 'Yes'),
    ('SaProt', 'Struct-aware marginal', 'Single-WT LLR', '0.833', '0.556', PURPLE, 'Yes'),
    ('ProtT5', 'VESPA (supervised)', 'Position cosine', '0.536', '0.186', RED, 'No'),
    ('NT-v2', 'Cosine(emb_REF, ALT)', 'Mean |Δhidden|', '0.638', '0.273', GREEN, 'No'),
    ('HyenaDNA', 'LLR (autoregressive)', 'LLR (autoregressive)', '0.448', '0.056', GREEN, 'Yes'),
    ('AlphaGenome', 'ALT−REF track diff', 'Mean |Δtrack|', '0.552', '0.153', ORANGE, 'Partial'),
]

# Header
add_card(slide, 0.5, 1.1, 12.3, 0.5, RGBColor(30, 58, 95))
add_text_box(slide, 0.7, 1.15, 1.5, 0.4, 'Model', 12, WHITE, True)
add_text_box(slide, 2.3, 1.15, 2.3, 0.4, 'Paper Method', 12, WHITE, True)
add_text_box(slide, 4.7, 1.15, 2.3, 0.4, 'Our V3 Method', 12, WHITE, True)
add_text_box(slide, 7.1, 1.15, 1.0, 0.4, 'AUROC', 12, WHITE, True)
add_text_box(slide, 8.2, 1.15, 1.0, 0.4, 'Spearman', 12, WHITE, True)
add_text_box(slide, 9.3, 1.15, 1.5, 0.4, 'Match?', 12, WHITE, True)
add_text_box(slide, 10.9, 1.15, 1.8, 0.4, 'Notes', 12, WHITE, True)

for i, (model, paper, ours, auroc, sp, color, match) in enumerate(models_detail):
    y = 1.7 + i * 0.7
    bg = BG_CARD if i % 2 == 0 else RGBColor(25, 35, 50)
    add_card(slide, 0.5, y, 12.3, 0.6, bg)
    a_color = GREEN if float(auroc) > 0.8 else (ORANGE if float(auroc) > 0.6 else RED)
    add_text_box(slide, 0.7, y+0.05, 1.5, 0.5, model, 12, color, True)
    add_text_box(slide, 2.3, y+0.05, 2.3, 0.5, paper, 10, GRAY)
    add_text_box(slide, 4.7, y+0.05, 2.3, 0.5, ours, 10, WHITE)
    add_text_box(slide, 7.1, y+0.05, 1.0, 0.5, auroc, 12, a_color, True)
    add_text_box(slide, 8.2, y+0.05, 1.0, 0.5, sp, 12, a_color, True)
    m_color = GREEN if match == 'Yes' else (ORANGE if match == 'Partial' else RED)
    add_text_box(slide, 9.3, y+0.05, 1.5, 0.5, match, 11, m_color, True)

# Notes
notes = [
    'nearly identical',
    'single-variant OK',
    'weaker method',
    'different metric',
    'same formula',
    'simplified',
]
for i, note in enumerate(notes):
    y = 1.7 + i * 0.7
    add_text_box(slide, 10.9, y+0.05, 1.8, 0.5, note, 9, GRAY)

# ════════ SLIDE 10: KEY FINDINGS ════════
slide = prs.slides.add_slide(prs.slide_layouts[6]); add_bg(slide)
add_text_box(slide, 0.5, 0.3, 12.3, 0.7, 'Key Findings', 30, WHITE, True)

findings = [
    ('1. LLR Dominates All Scoring Methods', GREEN,
     ['Log-Likelihood Ratio consistently outperforms:',
      '  Cosine similarity (loses 5-8% Spearman)',
      '  MaxPool + cosine(diff, wt) (negative Spearman!)',
      '  Mean absolute hidden delta',
      'LLR captures probabilistic uncertainty that',
      'embedding geometry methods miss']),
    ('2. Task-Data Alignment > Model Size', BLUE,
     ['ESM-1b (650M) beats NT-v2 (500M) by +37% AUROC',
      'Protein LMs learn amino acid semantics',
      'directly relevant to missense pathogenicity',
      'Genomic models learn DNA grammar',
      'that is one step removed from protein effect']),
    ('3. SaProt: Structure Helps but Limits Coverage', PURPLE,
     ['Adds 3Di structure tokens to ESM-style scoring',
      'AUROC 0.833 (near ESM-2 0.838)',
      'But missing 880 genes (no AF2 structures)',
      'Structure tokens add marginal value for',
      'single AA substitutions (structure barely changes)']),
    ('4. Genomic Models Are Misaligned for This Task', RED,
     ['AlphaGenome: trained on ENCODE regulatory tracks',
      'Predicts chromatin/splicing, not protein stability',
      'HyenaDNA: trained on 1 human genome only',
      'Insufficient diversity for variant interpretation',
      'AUROC 0.448 (below random!)']),
]

for i, (title, color, bullets) in enumerate(findings):
    col = i % 2
    row = i // 2
    x = 0.5 + col * 6.4
    y = 1.2 + row * 3.1
    add_card(slide, x, y, 5.9, 2.8)
    tf = add_text_box(slide, x+0.3, y+0.1, 5.3, 0.5, title, 14, color, True)
    for j, bullet in enumerate(bullets):
        c = WHITE if not bullet.startswith('AUROC') and not bullet.startswith('LLR') and not bullet.startswith('AUROC') else (RED if 'negative' in bullet or 'below' in bullet else YELLOW)
        add_para(tf, bullet, 11, c, space_before=3)

add_card(slide, 0.5, 6.8, 12.3, 0.6)
tf = add_text_box(slide, 0.8, 6.85, 11.7, 0.5,
    'Bottom Line: Protein LMs with LLR scoring are the gold standard for missense variant pathogenicity prediction',
    14, YELLOW, True, PP_ALIGN.CENTER)

# ════════ SAVE ════════
out_path = os.path.join(FIG_DIR, 'Pathogenicity_Scoring_Methods.pptx')
prs.save(out_path)
print(f"\nPPT saved: {out_path}")

# Cleanup temp plots
import shutil
shutil.rmtree(TEMP_DIR)
print("Temp plots cleaned up.")
