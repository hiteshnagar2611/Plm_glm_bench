#!/usr/bin/env python3
"""
Pathogenicity Scoring Methods PPT:
- How each model calculates pathogenicity (paper method)
- How we calculated it in V3 benchmark
- Results comparison
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)

# ── Color Palette ──
BG_DARK   = RGBColor(15, 23, 42)
BG_CARD   = RGBColor(30, 41, 59)
WHITE     = RGBColor(255, 255, 255)
GRAY      = RGBColor(148, 163, 184)
BLUE      = RGBColor(59, 130, 246)
GREEN     = RGBColor(16, 185, 129)
RED       = RGBColor(239, 68, 68)
ORANGE    = RGBColor(249, 115, 22)
PURPLE    = RGBColor(139, 92, 246)
CYAN      = RGBColor(6, 182, 212)
YELLOW    = RGBColor(234, 179, 8)

def add_bg(slide):
    bg = slide.background.fill
    bg.solid()
    bg.fore_color.rgb = BG_DARK

def add_text_box(slide, left, top, width, height, text, font_size=14, color=WHITE,
                 bold=False, alignment=PP_ALIGN.LEFT, font_name='Calibri'):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    return tf

def add_para(tf, text, font_size=14, color=WHITE, bold=False, space_before=6, alignment=PP_ALIGN.LEFT):
    p = tf.add_paragraph()
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = 'Calibri'
    p.space_before = Pt(space_before)
    p.alignment = alignment
    return p

def add_card(slide, left, top, width, height, fill_color=BG_CARD):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(height))
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()
    shape.shadow.inherit = False
    return shape

# ════════════════════════════════════════════════════════════
# SLIDE 1: TITLE
# ════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)

# Accent line
shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(1.5), Inches(2.6), Inches(10.3), Inches(0.06))
shape.fill.solid(); shape.fill.fore_color.rgb = BLUE; shape.line.fill.background()

add_text_box(slide, 1.5, 1.0, 10.3, 1.6,
             'Pathogenicity Scoring Methods', 42, WHITE, True, PP_ALIGN.CENTER)
add_text_box(slide, 1.5, 2.8, 10.3, 1.0,
             'How PLMs and GLMs Calculate Variant Pathogenicity', 22, GRAY, False, PP_ALIGN.CENTER)
add_text_box(slide, 1.5, 4.0, 10.3, 0.8,
             'Paper Methods  |  Our V3 Benchmark  |  Results Comparison', 16, BLUE, False, PP_ALIGN.CENTER)

# Bottom labels
add_text_box(slide, 1.5, 5.5, 5, 0.5, 'Protein Language Models', 14, BLUE, True, PP_ALIGN.LEFT)
add_text_box(slide, 6.8, 5.5, 5, 0.5, 'Genomic Language Models', 14, GREEN, True, PP_ALIGN.LEFT)
tf = add_text_box(slide, 1.5, 5.9, 5, 0.8, 'ESM-1b  |  ESM-2  |  SaProt  |  ProtT5', 12, GRAY, False, PP_ALIGN.LEFT)
tf2 = add_text_box(slide, 6.8, 5.9, 5, 0.8, 'NT-v2  |  HyenaDNA  |  AlphaGenome', 12, GRAY, False, PP_ALIGN.LEFT)

# ════════════════════════════════════════════════════════════
# SLIDE 2: WHAT IS PATHOGENICITY SCORING?
# ════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_text_box(slide, 0.5, 0.3, 12.3, 0.7, 'What Is Pathogenicity Scoring?', 30, WHITE, True)

# Problem definition card
add_card(slide, 0.5, 1.2, 12.3, 1.6)
tf = add_text_box(slide, 0.8, 1.35, 11.7, 1.4, 'The Problem', 18, BLUE, True)
add_para(tf, 'Given a genetic variant (single amino acid or nucleotide change), predict whether it is:', 14, WHITE, space_before=8)
add_para(tf, '   Pathogenic (disease-causing)  vs  Benign (harmless)', 16, GREEN, True, space_before=6)

# Two paradigms
add_card(slide, 0.5, 3.1, 5.9, 3.8)
tf = add_text_box(slide, 0.8, 3.2, 5.5, 0.5, 'Zero-Shot Scoring', 20, GREEN, True)
add_para(tf, 'No task-specific training needed', 13, GRAY, space_before=8)
add_para(tf, '', 6, GRAY, space_before=4)
add_para(tf, 'Protein LMs (ESM, SaProt, ProtT5):', 13, BLUE, True)
add_para(tf, '  Compare P(mutant) vs P(wildtype) from the', 13, WHITE)
add_para(tf, '  pretrained masked language model', 13, WHITE)
add_para(tf, '', 6, GRAY, space_before=4)
add_para(tf, 'Genomic LMs (NT-v2, HyenaDNA):', 13, GREEN, True)
add_para(tf, '  Compare embeddings or log-probabilities', 13, WHITE)
add_para(tf, '  for reference vs alternate allele', 13, WHITE)
add_para(tf, '', 6, GRAY, space_before=4)
add_para(tf, 'AlphaGenome:', 13, ORANGE, True)
add_para(tf, '  Compare predicted regulatory tracks', 13, WHITE)
add_para(tf, '  for REF vs ALT sequences', 13, WHITE)

add_card(slide, 6.9, 3.1, 5.9, 3.8)
tf = add_text_box(slide, 7.2, 3.2, 5.5, 0.5, 'Fine-Tuned Scoring', 20, ORANGE, True)
add_para(tf, 'Requires task-specific labeled data', 13, GRAY, space_before=8)
add_para(tf, '', 6, GRAY, space_before=4)
add_para(tf, 'How it works:', 13, WHITE, True)
add_para(tf, '  1. Pretrain on large unlabeled corpus', 13, WHITE)
add_para(tf, '  2. Fine-tune on labeled variant data', 13, WHITE)
add_para(tf, '     (e.g., ClinVar, ProteinGym DMS)', 13, WHITE)
add_para(tf, '  3. Train a classification head on top', 13, WHITE)
add_para(tf, '', 6, GRAY, space_before=4)
add_para(tf, 'Examples:', 13, WHITE, True)
add_para(tf, '  ESM-1v (variant-specific fine-tuning)', 13, GRAY)
add_para(tf, '  NT-v2 (supervised downstream head)', 13, GRAY)
add_para(tf, '  AlphaGenome (pretrained on ENCODE)', 13, GRAY)
add_para(tf, '', 6, GRAY, space_before=4)
add_para(tf, 'Our V3 benchmark uses zero-shot for all', 13, YELLOW, True)
add_para(tf, 'models to ensure fair comparison', 13, YELLOW)

# ════════════════════════════════════════════════════════════
# SLIDE 3: LOG-LIKELIHOOD RATIO (CORE METHOD)
# ════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_text_box(slide, 0.5, 0.3, 12.3, 0.7, 'Core Method: Log-Likelihood Ratio (LLR)', 30, WHITE, True)

# Main formula card
add_card(slide, 0.5, 1.2, 12.3, 2.0)
tf = add_text_box(slide, 0.8, 1.3, 11.7, 0.5, 'The Fundamental Formula', 18, BLUE, True)
add_para(tf, '', 4, GRAY, space_before=4)
add_para(tf, 'LLR  =  log P(mutant residue  |  mutant sequence context)', 20, GREEN, True, space_before=6)
add_para(tf, '            −  log P(wildtype residue  |  wildtype sequence context)', 20, GREEN, True, space_before=2)
add_para(tf, '', 4, GRAY, space_before=4)
add_para(tf, 'Negative score → model assigns lower probability to mutant → likely pathogenic', 14, GRAY, space_before=6)

# How it works step by step
add_card(slide, 0.5, 3.5, 5.9, 3.5)
tf = add_text_box(slide, 0.8, 3.6, 5.5, 0.5, 'Step-by-Step', 18, GREEN, True)
add_para(tf, '', 4, GRAY, space_before=2)
add_para(tf, '1. Input wildtype sequence → run through LM', 13, WHITE, space_before=6)
add_para(tf, '   Extract log P(wt_aa | context) at variant pos', 12, GRAY)
add_para(tf, '', 4, GRAY, space_before=2)
add_para(tf, '2. Input mutant sequence → run through LM', 13, WHITE, space_before=6)
add_para(tf, '   Extract log P(mut_aa | context) at variant pos', 12, GRAY)
add_para(tf, '', 4, GRAY, space_before=2)
add_para(tf, '3. Score = log P(mut) − log P(wt)', 13, WHITE, space_before=6)
add_para(tf, '', 4, GRAY, space_before=2)
add_para(tf, 'More negative → more pathogenic', 13, RED, True, space_before=6)
add_para(tf, 'Score ≥ 0 → model thinks mutation is neutral/beneficial', 13, GREEN)

# Visual example
add_card(slide, 6.9, 3.5, 5.9, 3.5)
tf = add_text_box(slide, 7.2, 3.6, 5.5, 0.5, 'Example: ESM-2 at Position 42', 18, BLUE, True)
add_para(tf, '', 4, GRAY, space_before=2)
add_para(tf, 'Wildtype: ...A R K [G] L E D...', 13, WHITE, space_before=8)
add_para(tf, '  P(Gly | context) = 0.85', 13, GRAY)
add_para(tf, '  log P(Gly) = −0.163', 13, GRAY)
add_para(tf, '', 4, GRAY, space_before=2)
add_para(tf, 'Mutant:   ...A R K [D] L E D...', 13, WHITE, space_before=8)
add_para(tf, '  P(Asp | context) = 0.02', 13, GRAY)
add_para(tf, '  log P(Asp) = −3.912', 13, GRAY)
add_para(tf, '', 4, GRAY, space_before=2)
add_para(tf, 'LLR = −3.912 − (−0.163) = −3.749', 14, RED, True, space_before=8)
add_para(tf, '→ Strongly pathogenic (charge reversal G→D)', 12, GRAY)

# ════════════════════════════════════════════════════════════
# SLIDE 4: ESM-1b & ESM-2 (PAPER METHOD)
# ════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_text_box(slide, 0.5, 0.3, 12.3, 0.7, 'ESM-1b & ESM-2: Paper Method', 30, WHITE, True)

# Paper method
add_card(slide, 0.5, 1.2, 12.3, 2.8)
tf = add_text_box(slide, 0.8, 1.3, 11.7, 0.5, 'Original Paper: Masked Marginal LLR (Meier et al. 2021, Lin et al. 2023)', 16, BLUE, True)
add_para(tf, '', 4, GRAY, space_before=2)
add_para(tf, 'Paper Formula:', 14, WHITE, True, space_before=6)
add_para(tf, '  Score = Σ [ log P(x_t = x_mut | x_masked) − log P(x_t = x_wt | x_masked) ]', 18, GREEN, True, space_before=6)
add_para(tf, '', 4, GRAY, space_before=2)
add_para(tf, 'Key difference from our method:', 14, ORANGE, True, space_before=6)
add_para(tf, '  Paper: Places [MASK] at ALL mutated positions simultaneously → mutual information', 13, WHITE)
add_para(tf, '  Ours:  Runs WT and mutant sequences independently → pseudolikelihood (no masking)', 13, WHITE)
add_para(tf, '', 4, GRAY, space_before=2)
add_para(tf, '  The paper method prevents information leakage between multiple mutations (important for DMS).', 12, GRAY)
add_para(tf, '  For single-variant ClinVar scoring, both methods give nearly identical results.', 12, GRAY)

# Our method
add_card(slide, 0.5, 4.3, 5.9, 2.8)
tf = add_text_box(slide, 0.8, 4.4, 5.5, 0.5, 'Our V3 Method', 16, GREEN, True)
add_para(tf, '', 4, GRAY, space_before=2)
add_para(tf, 'Model: esm1b_t33_650M_UR50S / esm2_t33_650M_UR50D', 12, GRAY, space_before=6)
add_para(tf, '', 4, GRAY, space_before=2)
add_para(tf, '1. Run WT sequence → logits at variant pos', 13, WHITE, space_before=6)
add_para(tf, '   wt_ll = log softmax(logits)[wt_token]', 12, GRAY)
add_para(tf, '', 4, GRAY, space_before=2)
add_para(tf, '2. Run mutant sequence → logits at variant pos', 13, WHITE, space_before=6)
add_para(tf, '   mut_ll = log softmax(logits)[mut_token]', 12, GRAY)
add_para(tf, '', 4, GRAY, space_before=2)
add_para(tf, '3. LLR = mut_ll − wt_ll', 13, GREEN, True, space_before=6)

# Results
add_card(slide, 6.9, 4.3, 5.9, 2.8)
tf = add_text_box(slide, 7.2, 4.4, 5.5, 0.5, 'V3 Results', 16, BLUE, True)
add_para(tf, '', 4, GRAY, space_before=2)
add_para(tf, 'ESM-1b (650M params, UniRef50):', 14, WHITE, True, space_before=8)
add_para(tf, '   AUROC = 0.873  |  Spearman = 0.631', 14, BLUE, True)
add_para(tf, '', 4, GRAY, space_before=4)
add_para(tf, 'ESM-2 (650M params, UniRef50):', 14, WHITE, True, space_before=8)
add_para(tf, '   AUROC = 0.838  |  Spearman = 0.572', 14, BLUE, True)
add_para(tf, '', 4, GRAY, space_before=4)
add_para(tf, 'ESM-1b paper ClinVar AUROC ≈ 0.89 (similar)', 12, GRAY, space_before=6)
add_para(tf, '→ Our results consistent with literature', 12, GREEN)

# ════════════════════════════════════════════════════════════
# SLIDE 5: SAPROT (STRUCTURE-AWARE)
# ════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_text_box(slide, 0.5, 0.3, 12.3, 0.7, 'SaProt: Structure-Aware Scoring', 30, WHITE, True)

# Paper method
add_card(slide, 0.5, 1.2, 12.3, 3.0)
tf = add_text_box(slide, 0.8, 1.3, 11.7, 0.5, 'Paper: Structure-Aware Masked Marginal LLR (Su et al. 2024, ICLR)', 16, BLUE, True)
add_para(tf, '', 4, GRAY, space_before=2)
add_para(tf, 'Innovation: Combines amino acid + 3Di structural token from Foldseek', 14, WHITE, space_before=6)
add_para(tf, '', 4, GRAY, space_before=2)
add_para(tf, 'Paper Formula:', 14, WHITE, True, space_before=6)
add_para(tf, '  Score = Σ_f [ log P(mut_f | context) − log P(wt_f | context) ]', 18, GREEN, True, space_before=6)
add_para(tf, '  where f iterates over all Foldseek structural tokens for each residue type', 12, GRAY, space_before=4)
add_para(tf, '', 4, GRAY, space_before=2)
add_para(tf, 'Vocabulary:  AA × 3Di  →  e.g., "Da" = Aspartate in helix, "Kc" = Lysine in coil', 13, WHITE, space_before=6)
add_para(tf, '', 4, GRAY, space_before=2)
add_para(tf, 'Paper uses structures from AlphaFoldDB (~40M proteins with predicted structures)', 13, GRAY, space_before=6)

# Our method
add_card(slide, 0.5, 4.5, 5.9, 2.6)
tf = add_text_box(slide, 0.8, 4.6, 5.5, 0.5, 'Our V3 Method', 16, GREEN, True)
add_para(tf, '', 4, GRAY, space_before=2)
add_para(tf, 'Model: westlake-repl/SaProt_650M_AF2', 12, GRAY, space_before=6)
add_para(tf, 'Uses only WT sequence (masked-LM style):', 13, WHITE, space_before=6)
add_para(tf, '  wt_token = "Ga" (Glycine in helix)', 12, GRAY)
add_para(tf, '  mut_token = "Da" (Aspartate in helix)', 12, GRAY)
add_para(tf, '', 4, GRAY, space_before=2)
add_para(tf, '  Score = log P(mut_token | wt_context)', 13, GREEN, True, space_before=6)
add_para(tf, '             − log P(wt_token | wt_context)', 13, GREEN, True)
add_para(tf, '', 4, GRAY, space_before=2)
add_para(tf, 'Only WT run needed (mutant scored from same context)', 12, ORANGE, space_before=6)

# Results
add_card(slide, 6.9, 4.5, 5.9, 2.6)
tf = add_text_box(slide, 7.2, 4.6, 5.5, 0.5, 'V3 Results', 16, BLUE, True)
add_para(tf, '', 4, GRAY, space_before=2)
add_para(tf, 'SaProt (650M params, UniRef50 + AF2):', 14, WHITE, True, space_before=8)
add_para(tf, '   AUROC = 0.833  |  Spearman = 0.556', 14, BLUE, True)
add_para(tf, '', 4, GRAY, space_before=4)
add_para(tf, 'Paper ProteinGym Spearman ≈ 0.45 (DMS)', 12, GRAY, space_before=6)
add_para(tf, '→ Comparable; SaProt missing 880 genes (no AF2 structures)', 12, ORANGE, space_before=4)

# ════════════════════════════════════════════════════════════
# SLIDE 6: PROTT5 (EMBEDDING METHOD)
# ════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_text_box(slide, 0.5, 0.3, 12.3, 0.7, 'ProtT5: Embedding-Based Scoring', 30, WHITE, True)

# Paper method
add_card(slide, 0.5, 1.2, 12.3, 2.5)
tf = add_text_box(slide, 0.8, 1.3, 11.7, 0.5, 'Paper: VESPA Ensemble (Marquet et al. 2021) or Direct LLR (Elnaggar et al. 2022)', 16, BLUE, True)
add_para(tf, '', 4, GRAY, space_before=2)
add_para(tf, 'VESPA (supervised): Combines ProtT5 embeddings + BLOSUM62 + conservation → logistic regression', 13, WHITE, space_before=6)
add_para(tf, 'Direct LLR: Same masked marginal formula as ESM-1b (log P(mut) − log P(wt))', 13, WHITE, space_before=4)
add_para(tf, '', 4, GRAY, space_before=2)
add_para(tf, 'Architecture: T5 encoder-decoder (3B params), only encoder used for embeddings', 13, GRAY, space_before=6)

# Our method
add_card(slide, 0.5, 4.0, 5.9, 3.1)
tf = add_text_box(slide, 0.8, 4.1, 5.5, 0.5, 'Our V3 Method', 16, GREEN, True)
add_para(tf, '', 4, GRAY, space_before=2)
add_para(tf, 'Model: Rostlab/prot_t5_xl_uniref50', 12, GRAY, space_before=6)
add_para(tf, 'Uses position-level cosine similarity:', 13, WHITE, space_before=6)
add_para(tf, '  wt_emb  = encoder(sequence_wt)[pos]', 12, GRAY)
add_para(tf, '  mut_emb = encoder(sequence_mut)[pos]', 12, GRAY)
add_para(tf, '', 4, GRAY, space_before=2)
add_para(tf, '  Score = cosine(wt_emb, mut_emb)', 13, GREEN, True, space_before=6)
add_para(tf, '', 4, GRAY, space_before=2)
add_para(tf, 'Lower cosine → less similar → more pathogenic', 13, ORANGE, space_before=6)
add_para(tf, '', 4, GRAY, space_before=2)
add_para(tf, '(Why not LLR? ProtT5 is encoder-decoder, not a', 12, GRAY, space_before=4)
add_para(tf, ' standard MLM like ESM. Position cosine is more', 12, GRAY)
add_para(tf, ' natural for T5 encoder embeddings.)', 12, GRAY)

# Results
add_card(slide, 6.9, 4.0, 5.9, 3.1)
tf = add_text_box(slide, 7.2, 4.1, 5.5, 0.5, 'V3 Results', 16, BLUE, True)
add_para(tf, '', 4, GRAY, space_before=2)
add_para(tf, 'ProtT5-XL (3B params, UniRef50):', 14, WHITE, True, space_before=8)
add_para(tf, '   AUROC = 0.536  |  Spearman = 0.186', 14, RED, True)
add_para(tf, '', 4, GRAY, space_before=4)
add_para(tf, 'Why poor performance?', 14, ORANGE, True, space_before=8)
add_para(tf, '  1. Position cosine ≠ LLR (loses probability info)', 12, WHITE, space_before=6)
add_para(tf, '  2. ProtT5 encoder not optimized for', 12, WHITE, space_before=4)
add_para(tf, '     single-position variant scoring', 12, WHITE)
add_para(tf, '  3. VESPA (supervised) would perform better', 12, WHITE, space_before=4)
add_para(tf, '     but requires labeled training data', 12, WHITE)
add_para(tf, '', 4, GRAY, space_before=4)
add_para(tf, '→ Zero-shot cosine is a weak scoring method for ProtT5', 12, YELLOW)

# ════════════════════════════════════════════════════════════
# SLIDE 7: NT-v2 (GENOMIC)
# ════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_text_box(slide, 0.5, 0.3, 12.3, 0.7, 'Nucleotide Transformer v2: Genomic Scoring', 30, WHITE, True)

# Paper method
add_card(slide, 0.5, 1.2, 12.3, 2.8)
tf = add_text_box(slide, 0.8, 1.3, 11.7, 0.5, 'Paper: Multiple Zero-Shot Metrics (Dalla-Torre et al. 2024, Nature Methods)', 16, BLUE, True)
add_para(tf, '', 4, GRAY, space_before=2)
add_para(tf, 'Paper tested several zero-shot scoring methods:', 14, WHITE, space_before=6)
add_para(tf, '', 4, GRAY, space_before=2)
add_para(tf, '  Best:  Cosine similarity = cos(emb_REF, emb_ALT)', 14, GREEN, True, space_before=6)
add_para(tf, '  Also:  Dot product, Manhattan distance, Euclidean distance', 13, WHITE, space_before=4)
add_para(tf, '         Jensen-Shannon divergence, Hellinger distance', 13, WHITE, space_before=4)
add_para(tf, '', 4, GRAY, space_before=2)
add_para(tf, 'Input: 6kb (v1) or 12kb (v2) window centered on variant', 13, GRAY, space_before=6)
add_para(tf, 'Lower cosine → more pathogenic (paper)', 13, ORANGE, space_before=4)

# Our method
add_card(slide, 0.5, 4.3, 5.9, 2.8)
tf = add_text_box(slide, 0.8, 4.4, 5.5, 0.5, 'Our V3 Method', 16, GREEN, True)
add_para(tf, '', 4, GRAY, space_before=2)
add_para(tf, 'Model: nucleotide-transformer-v2-500m-multi-species', 11, GRAY, space_before=6)
add_para(tf, 'Uses mean absolute hidden-state delta:', 13, WHITE, space_before=6)
add_para(tf, '  Input: 6001bp DNA (variant at pos 3000)', 12, GRAY)
add_para(tf, '  Run WT sequence → hidden_states_wt', 12, GRAY)
add_para(tf, '  Run mutant sequence → hidden_states_mut', 12, GRAY)
add_para(tf, '', 4, GRAY, space_before=2)
add_para(tf, '  Score = mean(|hidden_wt − hidden_mut|)', 13, GREEN, True, space_before=6)
add_para(tf, '', 4, GRAY, space_before=2)
add_para(tf, 'Higher delta → more embedding shift → more pathogenic', 13, ORANGE, space_before=6)
add_para(tf, '', 4, GRAY, space_before=2)
add_para(tf, '(Why not cosine? Hidden-state delta captures overall', 11, GRAY, space_before=4)
add_para(tf, ' embedding magnitude change, not just direction.)', 11, GRAY)

# Results
add_card(slide, 6.9, 4.3, 5.9, 2.8)
tf = add_text_box(slide, 7.2, 4.4, 5.5, 0.5, 'V3 Results', 16, BLUE, True)
add_para(tf, '', 4, GRAY, space_before=2)
add_para(tf, 'NT-v2 (500M params, 850 genomes):', 14, WHITE, True, space_before=8)
add_para(tf, '   AUROC = 0.638  |  Spearman = 0.273', 14, BLUE, True)
add_para(tf, '', 4, GRAY, space_before=4)
add_para(tf, 'Why moderate performance?', 14, ORANGE, True, space_before=8)
add_para(tf, '  1. Genomic model on protein-coding variants', 12, WHITE, space_before=6)
add_para(tf, '     (trained on DNA, not protein sequences)', 12, WHITE)
add_para(tf, '  2. 6kb window truncates to 2048 tokens', 12, WHITE, space_before=4)
add_para(tf, '  3. NT-v2 is designed for regulatory tasks', 12, WHITE, space_before=4)
add_para(tf, '     not protein pathogenicity', 12, WHITE)
add_para(tf, '', 4, GRAY, space_before=4)
add_para(tf, '→ Better than random but below protein LMs', 12, YELLOW)

# ════════════════════════════════════════════════════════════
# SLIDE 8: HYENADNA & ALPHAGENOME
# ════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_text_box(slide, 0.5, 0.3, 12.3, 0.7, 'HyenaDNA & AlphaGenome: Genomic Scoring', 30, WHITE, True)

# HyenaDNA card
add_card(slide, 0.5, 1.2, 5.9, 5.8)
tf = add_text_box(slide, 0.8, 1.3, 5.5, 0.5, 'HyenaDNA (Nguyen et al. 2023)', 18, GREEN, True)
add_para(tf, '', 4, GRAY, space_before=2)
add_para(tf, 'Paper Method:', 14, WHITE, True, space_before=6)
add_para(tf, '  Autoregressive LLR:', 13, GREEN, True, space_before=6)
add_para(tf, '  Score = log P(seq_ALT) − log P(seq_REF)', 14, GREEN, True, space_before=4)
add_para(tf, '', 4, GRAY, space_before=2)
add_para(tf, '  HyenaDNA is decoder-only (causal LM)', 12, GRAY, space_before=6)
add_para(tf, '  Context flows left→right only', 12, GRAY, space_before=4)
add_para(tf, '  Single nucleotide tokens (no k-mer)', 12, GRAY, space_before=4)
add_para(tf, '', 4, GRAY, space_before=2)
add_para(tf, 'Our V3 Method:', 14, WHITE, True, space_before=8)
add_para(tf, '  Same LLR formula (paper method):', 13, WHITE, space_before=6)
add_para(tf, '  ref_ll = log P(ref_base | ref_context)', 12, GRAY)
add_para(tf, '  alt_ll = log P(alt_base | alt_context)', 12, GRAY)
add_para(tf, '  LLR = alt_ll − ref_ll', 13, GREEN, True, space_before=4)
add_para(tf, '', 4, GRAY, space_before=2)
add_para(tf, 'V3 Results:', 14, WHITE, True, space_before=8)
add_para(tf, '  HyenaDNA (150M, hg38):', 13, WHITE, space_before=6)
add_para(tf, '  AUROC = 0.448  |  Spearman = 0.056', 14, RED, True)
add_para(tf, '', 4, GRAY, space_before=2)
add_para(tf, '  → Worst performer (1 genome training data)', 12, YELLOW, space_before=6)

# AlphaGenome card
add_card(slide, 6.9, 1.2, 5.9, 5.8)
tf = add_text_box(slide, 7.2, 1.3, 5.5, 0.5, 'AlphaGenome (Avsec et al. 2025)', 18, ORANGE, True)
add_para(tf, '', 4, GRAY, space_before=2)
add_para(tf, 'Paper Method:', 14, WHITE, True, space_before=6)
add_para(tf, '  Multi-track difference:', 13, ORANGE, True, space_before=6)
add_para(tf, '  Score = Aggregate(ALT tracks) − Aggregate(REF tracks)', 13, ORANGE, True, space_before=4)
add_para(tf, '', 4, GRAY, space_before=2)
add_para(tf, '  Per modality (paper):', 12, GRAY, space_before=6)
add_para(tf, '  RNA-seq: log(mean(ALT)/mean(REF))', 12, WHITE, space_before=4)
add_para(tf, '  Splice: max(|ALT − REF|) class probs', 12, WHITE, space_before=4)
add_para(tf, '  DNase: log2((sum(ALT)+1)/(sum(REF)+1))', 12, WHITE, space_before=4)
add_para(tf, '', 4, GRAY, space_before=2)
add_para(tf, 'Our V3 Method:', 14, WHITE, True, space_before=8)
add_para(tf, '  Mean absolute delta across 5 tracks:', 13, WHITE, space_before=6)
add_para(tf, '  For each output type i:', 12, GRAY)
add_para(tf, '    delta_i = mean(|ref_i − alt_i|)', 12, GRAY)
add_para(tf, '  Score = mean(deltas across all tracks)', 13, GREEN, True, space_before=4)
add_para(tf, '', 4, GRAY, space_before=2)
add_para(tf, 'V3 Results:', 14, WHITE, True, space_before=8)
add_para(tf, '  AlphaGenome (250M, ENCODE+GTEx):', 13, WHITE, space_before=6)
add_para(tf, '  AUROC = 0.552  |  Spearman = 0.153', 14, ORANGE, True)
add_para(tf, '', 4, GRAY, space_before=2)
add_para(tf, '  → Regulatory model, not designed for', 12, YELLOW, space_before=6)
add_para(tf, '    protein-coding missense pathogenicity', 12, YELLOW)

# ════════════════════════════════════════════════════════════
# SLIDE 9: COMPARISON TABLE
# ════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_text_box(slide, 0.5, 0.3, 12.3, 0.7, 'Method Comparison: Paper vs Our V3', 30, WHITE, True)

# Table header
add_card(slide, 0.5, 1.2, 12.3, 0.6, RGBColor(30, 58, 95))
tf = add_text_box(slide, 0.7, 1.25, 1.8, 0.5, 'Model', 13, WHITE, True)
add_text_box(slide, 2.5, 1.25, 2.0, 0.5, 'Paper Method', 13, WHITE, True)
add_text_box(slide, 4.6, 1.25, 2.5, 0.5, 'Our V3 Method', 13, WHITE, True)
add_text_box(slide, 7.2, 1.25, 1.3, 0.5, 'AUROC', 13, WHITE, True)
add_text_box(slide, 8.6, 1.25, 1.3, 0.5, 'Spearman', 13, WHITE, True)
add_text_box(slide, 10.0, 1.25, 2.8, 0.5, 'Match?', 13, WHITE, True)

# Table rows
rows = [
    ('ESM-1b',   'Masked marginal LLR',        'Pseudolikelihood LLR',   '0.873', '0.631', 'Yes (nearly identical)'),
    ('ESM-2',    'Masked marginal LLR',        'Pseudolikelihood LLR',   '0.838', '0.572', 'Yes (nearly identical)'),
    ('SaProt',   'Structure-aware marginal LLR','Single-WT LLR',         '0.833', '0.556', 'Yes (single-variant)'),
    ('ProtT5',   'VESPA (supervised) / LLR',   'Position cosine',        '0.536', '0.186', 'No (weaker method)'),
    ('NT-v2',    'Cosine(emb_REF, emb_ALT)',   'Mean |Δhidden|',         '0.638', '0.273', 'No (different metric)'),
    ('HyenaDNA', 'LLR (autoregressive)',        'LLR (autoregressive)',   '0.448', '0.056', 'Yes (same formula)'),
    ('AlphaGenome','ALT−REF track difference',  'Mean |Δtrack|',          '0.552', '0.153', 'Partial (simplified)'),
]

for i, (model, paper, ours, auroc, sp, match) in enumerate(rows):
    y = 1.9 + i * 0.7
    bg = BG_CARD if i % 2 == 0 else RGBColor(25, 35, 50)
    add_card(slide, 0.5, y, 12.3, 0.6, bg)
    color = BLUE if float(auroc) > 0.8 else (ORANGE if float(auroc) > 0.6 else RED)
    add_text_box(slide, 0.7, y+0.05, 1.8, 0.5, model, 12, WHITE, True)
    add_text_box(slide, 2.5, y+0.05, 2.0, 0.5, paper, 11, GRAY)
    add_text_box(slide, 4.6, y+0.05, 2.5, 0.5, ours, 11, WHITE)
    add_text_box(slide, 7.2, y+0.05, 1.3, 0.5, auroc, 13, color, True)
    add_text_box(slide, 8.6, y+0.05, 1.3, 0.5, sp, 13, color, True)
    add_text_box(slide, 10.0, y+0.05, 2.8, 0.5, match, 11, GREEN if 'Yes' in match else (ORANGE if 'Partial' in match else RED))

# Key insight
add_card(slide, 0.5, 6.8, 12.3, 0.6)
tf = add_text_box(slide, 0.8, 6.85, 11.7, 0.5,
    'Key: Protein LMs (LLR) dominate → AUROC 0.83-0.87  |  Genomic LMs struggle → AUROC 0.45-0.64  |  Task mismatch matters most',
    13, YELLOW, True, PP_ALIGN.CENTER)

# ════════════════════════════════════════════════════════════
# SLIDE 10: KEY FINDINGS
# ════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_text_box(slide, 0.5, 0.3, 12.3, 0.7, 'Key Findings', 30, WHITE, True)

# Finding 1
add_card(slide, 0.5, 1.2, 5.9, 2.5)
tf = add_text_box(slide, 0.8, 1.3, 5.5, 0.5, '1. LLR Dominates All Scoring Methods', 16, GREEN, True)
add_para(tf, '', 4, GRAY, space_before=2)
add_para(tf, 'Log-Likelihood Ratio consistently outperforms:', 13, WHITE, space_before=6)
add_para(tf, '  Cosine similarity (loses 5-8% Spearman)', 13, GRAY, space_before=4)
add_para(tf, '  Max-pool + cosine(diff, wt) (negative Spearman!)', 13, GRAY, space_before=4)
add_para(tf, '  Mean absolute hidden delta', 13, GRAY, space_before=4)
add_para(tf, '', 4, GRAY, space_before=2)
add_para(tf, 'LLR captures probabilistic uncertainty that', 12, YELLOW, space_before=4)
add_para(tf, 'embedding geometry methods miss', 12, YELLOW)

# Finding 2
add_card(slide, 6.9, 1.2, 5.9, 2.5)
tf = add_text_box(slide, 7.2, 1.3, 5.5, 0.5, '2. Task-Data Alignment > Model Size', 16, BLUE, True)
add_para(tf, '', 4, GRAY, space_before=2)
add_para(tf, 'ESM-1b (650M, UniRef50) beats NT-v2 (500M, 850 genomes):', 13, WHITE, space_before=6)
add_para(tf, '  0.873 vs 0.638 AUROC (+37%)', 13, BLUE, True, space_before=4)
add_para(tf, '', 4, GRAY, space_before=2)
add_para(tf, 'Why: Protein LMs learn amino acid semantics', 12, GRAY, space_before=4)
add_para(tf, 'directly relevant to missense pathogenicity', 12, GRAY)
add_para(tf, 'while genomic models learn DNA grammar', 12, GRAY, space_before=2)
add_para(tf, 'that is one step removed from protein effect', 12, GRAY)

# Finding 3
add_card(slide, 0.5, 4.0, 5.9, 2.5)
tf = add_text_box(slide, 0.8, 4.1, 5.5, 0.5, '3. SaProt: Structure Helps but Limits Coverage', 16, PURPLE, True)
add_para(tf, '', 4, GRAY, space_before=2)
add_para(tf, 'SaProt adds 3Di structure tokens to ESM-style scoring', 13, WHITE, space_before=6)
add_para(tf, '  → AUROC 0.833 (near ESM-2 0.838)', 13, GRAY, space_before=4)
add_para(tf, '  → But missing 880 genes (no AlphaFold structures)', 13, ORANGE, space_before=4)
add_para(tf, '  → Structure tokens add marginal value for single AA', 13, GRAY, space_before=4)
add_para(tf, '    substitutions (structure barely changes)', 13, GRAY)

# Finding 4
add_card(slide, 6.9, 4.0, 5.9, 2.5)
tf = add_text_box(slide, 7.2, 4.1, 5.5, 0.5, '4. Genomic Models Are Misaligned for This Task', 16, RED, True)
add_para(tf, '', 4, GRAY, space_before=2)
add_para(tf, 'AlphaGenome: Trained on ENCODE regulatory tracks', 13, WHITE, space_before=6)
add_para(tf, '  → Predicts chromatin/splicing, not protein stability', 13, GRAY, space_before=4)
add_para(tf, '  → AUROC 0.552 (near random for missense)', 13, GRAY, space_before=4)
add_para(tf, '', 4, GRAY, space_before=2)
add_para(tf, 'HyenaDNA: Trained on 1 human genome only', 13, WHITE, space_before=6)
add_para(tf, '  → Insufficient diversity for variant interpretation', 13, GRAY, space_before=4)
add_para(tf, '  → AUROC 0.448 (below random!)', 13, GRAY, space_before=4)

# Bottom summary
add_card(slide, 0.5, 6.8, 12.3, 0.6)
tf = add_text_box(slide, 0.8, 6.85, 11.7, 0.5,
    'Bottom Line: Protein LMs with LLR scoring are the gold standard for missense variant pathogenicity prediction',
    14, YELLOW, True, PP_ALIGN.CENTER)

# ════════════════════════════════════════════════════════════
# SAVE
# ════════════════════════════════════════════════════════════
out_dir = os.path.join(os.path.dirname(__file__), '..', 'benchmark_v3', 'figures')
os.makedirs(out_dir, exist_ok=True)
out_path = os.path.join(out_dir, 'Pathogenicity_Scoring_Methods.pptx')
prs.save(out_path)
print(f"Saved: {out_path}")
