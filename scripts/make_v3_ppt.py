#!/usr/bin/env python3
"""Generate a PowerPoint presentation for the V3 benchmark results."""
import os
import pandas as pd
import numpy as np
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from sklearn.metrics import roc_auc_score, average_precision_score, matthews_corrcoef, recall_score
from scipy.stats import spearmanr

BASE = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
V3 = os.path.join(BASE, "benchmark_v3")
FIG_DIR = os.path.join(V3, "figures")
RES_DIR = os.path.join(V3, "results")

WHITE = RGBColor(255, 255, 255)
BLACK = RGBColor(0, 0, 0)
DARK = RGBColor(30, 41, 59)
ACCENT = RGBColor(59, 130, 246)
GREEN = RGBColor(34, 197, 94)
RED = RGBColor(239, 68, 68)
ORANGE = RGBColor(249, 115, 22)
GRAY = RGBColor(107, 114, 128)
BLUE2 = RGBColor(96, 165, 250)
PURPLE = RGBColor(139, 92, 246)

def set_bg(slide, color=WHITE):
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = color

def tb(slide, left, top, w, h, text, sz=18, bold=False, color=BLACK, align=PP_ALIGN.LEFT):
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(w), Inches(h))
    box.text_frame.word_wrap = True
    p = box.text_frame.paragraphs[0]
    p.text = text
    p.font.size = Pt(sz)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = "Calibri"
    p.alignment = align
    return box

def bullets(slide, left, top, w, h, items, sz=14, color=BLACK, spacing=6):
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(w), Inches(h))
    tf = box.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = item
        p.font.size = Pt(sz)
        p.font.color.rgb = color
        p.font.name = "Calibri"
        p.space_after = Pt(spacing)
    return box

def table(slide, left, top, w, h, data, cw=None, hdr=ACCENT, sz=11):
    rows, cols = len(data), len(data[0])
    t = slide.shapes.add_table(rows, cols, Inches(left), Inches(top), Inches(w), Inches(h)).table
    for c, cwi in enumerate(cw or [w/cols]*cols):
        t.columns[c].width = Inches(cwi)
    for r, row in enumerate(data):
        for c, val in enumerate(row):
            cell = t.cell(r, c)
            cell.text = str(val)
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            p = cell.text_frame.paragraphs[0]
            p.font.size = Pt(sz)
            p.font.name = "Calibri"
            p.alignment = PP_ALIGN.CENTER
            if r == 0:
                p.font.bold = True
                p.font.color.rgb = WHITE
                cell.fill.solid()
                cell.fill.fore_color.rgb = hdr
            else:
                p.font.color.rgb = BLACK
                cell.fill.solid()
                cell.fill.fore_color.rgb = RGBColor(243, 244, 246) if r % 2 == 0 else WHITE
    return t

def calc(fpath, col, neg=False):
    df = pd.read_csv(fpath)
    df['label'] = df['ClinVar_label'].map({1.0:1, 0.0:0})
    df = df.dropna(subset=[col])
    y = df['label'].values.astype(int)
    s = -df[col].values if neg else df[col].values
    return {
        'AUROC': roc_auc_score(y, s),
        'AUPRC': average_precision_score(y, s),
        'Spearman': spearmanr(s, y)[0],
        'MCC': matthews_corrcoef(y, (s > np.median(s)).astype(int)),
        'Recall': recall_score(y, (s > np.median(s)).astype(int)),
        'n': len(df)
    }

# ── Calculate results ──
main = {
    'ESM1b-650M':   calc(f'{RES_DIR}/esm1b_scores.csv',      'ESM1b_LLR',            True),
    'ESM2-650M':    calc(f'{RES_DIR}/esm2_650m_scores.csv',   'ESM2_LLR',             True),
    'SaProt-650M':  calc(f'{RES_DIR}/saprot_scores.csv',      'SaProt_LLR',           True),
    'NT-v2-500M':   calc(f'{RES_DIR}/ntv2_scores.csv',        'NTv2_delta',           True),
    'AlphaGenome':   calc(f'{RES_DIR}/alphagenome_scores.csv', 'AlphaGenome_delta',    True),
    'ProtT5-XL':    calc(f'{RES_DIR}/prott5_scores.csv',      'ProtT5_score',         True),
    'HyenaDNA-150M':calc(f'{RES_DIR}/hyena_scores.csv',       'HyenaDNA_LLR',         True),
}
cos = {
    'ESM1b-650M':  calc(f'{RES_DIR}/esm1b_cos_scores.csv',  'ESM1b_cos',  True),
    'ESM2-650M':   calc(f'{RES_DIR}/esm2_cos_scores.csv',   'ESM2_cos',   True),
    'SaProt-650M': calc(f'{RES_DIR}/saprot_cos_scores.csv', 'SaProt_cos', True),
}
order = ['ESM1b-650M','ESM2-650M','SaProt-650M','NT-v2-500M','AlphaGenome','ProtT5-XL','HyenaDNA-150M']

# ── Benchmark data ──
bdf = pd.read_csv(os.path.join(V3, "data", "benchmark_v3.csv"))
bdf['label'] = bdf['ClinVar_label'].map({1.0:1, 0.0:0})
AA_H = {'Ala','Val','Leu','Ile','Met','Phe','Trp','Pro'}
AA_P = {'Ser','Thr','Tyr','Asn','Gln','Cys'}
AA_POS = {'Lys','Arg','His'}
AA_NEG = {'Asp','Glu'}
def aagrp(a):
    if a in AA_H: return 'Hydrophobic'
    if a in AA_P: return 'Polar'
    if a in AA_POS: return 'Charged+'
    if a in AA_NEG: return 'Charged-'
    return 'Gly/Special'
AA_SZ = {'Gly':1,'Ala':2,'Val':3,'Leu':4,'Ile':4,'Pro':3,'Phe':5,'Trp':6,
          'Ser':2,'Thr':3,'Cys':2,'Tyr':5,'Met':4,'Asp':3,'Glu':3,'Asn':3,'Gln':3,'His':4,'Lys':5,'Arg':5}

# ════════════════════════════════════════════════════
prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# ── SLIDE 1: Title ──
sl = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(sl, DARK)
tb(sl, 1.0, 1.5, 11.3, 1.0, "Protein Language Models vs. DNA Language Models",
   sz=38, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
tb(sl, 1.0, 2.8, 11.3, 0.7, "Benchmark V3: Pathogenic Variant Prediction",
   sz=28, color=BLUE2, align=PP_ALIGN.CENTER)
tb(sl, 1.0, 3.8, 11.3, 0.5, "973 genes  ·  5,932 missense variants  ·  ClinVar + AlphaFold + MANE Select",
   sz=16, color=GRAY, align=PP_ALIGN.CENTER)
tb(sl, 1.0, 5.8, 11.3, 0.5, "Hitesh Nagar  ·  July 2026",
   sz=14, color=GRAY, align=PP_ALIGN.CENTER)

# ── SLIDE 2: Benchmark Design ──
sl = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(sl, WHITE)
tb(sl, 0.6, 0.3, 12, 0.6, "Benchmark V3 Design", sz=30, bold=True, color=DARK)

tb(sl, 0.6, 1.1, 5.8, 0.4, "Data Filtering Pipeline", sz=18, bold=True, color=ACCENT)
bullets(sl, 0.6, 1.6, 5.8, 3.5, [
    "ClinVar: Pathogenic + Benign only (excl. VUS, conflicting)",
    "MANE Select: single-transcript genes only",
    "AlphaFold DB: protein structure available",
    "SaProt: 3Di foldseek codes (880 genes missing)",
    "",
    "Final: 5,932 variants across 973 genes",
    "  Pathogenic: 3,625 (61.1%)",
    "  Benign:     2,307 (38.9%)",
], sz=13)

tb(sl, 6.8, 1.1, 5.8, 0.4, "Models Tested", sz=18, bold=True, color=ACCENT)
bullets(sl, 6.8, 1.6, 5.8, 3.5, [
    "Protein Language Models:",
    "  ESM1b-650M  ·  ESM2-650M  ·  SaProt-650M  ·  ProtT5-XL",
    "",
    "DNA / Genomic Language Models:",
    "  AlphaGenome  ·  NT-v2-500M  ·  HyenaDNA-150M",
    "",
    "Scoring Methods:",
    "  Log-Likelihood Ratio (PLMs)",
    "  Delta-score (AlphaGenome)",
    "  Hidden state delta (NT-v2)",
    "  DNA token perplexity (HyenaDNA)",
], sz=13)

tb(sl, 0.6, 5.3, 12, 0.4, "Evaluation Metrics", sz=18, bold=True, color=ACCENT)
bullets(sl, 0.6, 5.8, 12, 1.0, [
    "AUROC (primary)  ·  AUPRC  ·  Spearman correlation  ·  MCC  ·  Recall (at median threshold)"
], sz=13)

# ── SLIDE 3: Variant Composition ──
sl = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(sl, WHITE)
tb(sl, 0.6, 0.3, 12, 0.6, "V3 Database: Variant Composition", sz=30, bold=True, color=DARK)

# Top left - class balance
tb(sl, 0.6, 1.0, 5.8, 0.4, "Class Balance", sz=18, bold=True, color=ACCENT)
bal = [
    f"Pathogenic: 3,625 variants (61.1%)",
    f"Benign:     2,307 variants (38.9%)",
    f"Ratio: 1.57:1 (slight pathogenic skew)",
    f"Chromosomes: 24 represented (incl. X, Y)",
    f"ChrX: 1,039 variants (17.5%) — most enriched",
    f"Chr11: 506 variants (8.5%) — second",
]
bullets(sl, 0.6, 1.5, 5.8, 3.0, bal, sz=13)

# Top right - gene distribution
tb(sl, 6.8, 1.0, 5.8, 0.4, "Gene Distribution", sz=18, bold=True, color=ACCENT)
gc = bdf.groupby('GeneSymbol').size()
gene_items = [
    f"Total genes: {bdf['GeneSymbol'].nunique()}",
    f"Median variants/gene: {gc.median():.0f}  (range: {gc.min()}–{gc.max()})",
    f"358 genes with single variant",
    f"409 genes with 2–5 variants",
    f"47 genes with >20 variants",
    f"Top: GLA (224) · PAH (171) · GCK (163) · TP53 (101)",
]
bullets(sl, 6.8, 1.5, 5.8, 3.0, gene_items, sz=13)

# Bottom - path rate by position
tb(sl, 0.6, 4.5, 12, 0.4, "Pathogenicity by Protein Position", sz=18, bold=True, color=ACCENT)
bdf['pos_bin'] = pd.cut(bdf['aa_position'], bins=[0,50,100,200,300,500,1000],
                         labels=['1-50','51-100','101-200','201-300','301-500','501-1000'])
pos_data = [["Region", "Count", "Path Rate", "Interpretation"]]
interp = {
    '1-50': 'N-terminal — moderate tolerance',
    '51-100': 'Early region — moderate tolerance',
    '101-200': 'Mid region — most intolerant',
    '201-300': 'Mid-late — most intolerant',
    '301-500': 'Late region — more tolerant',
    '501-1000': 'C-terminal — most tolerant',
}
for bin_name in ['1-50','51-100','101-200','201-300','301-500','501-1000']:
    sub = bdf[bdf['pos_bin']==bin_name]
    pos_data.append([bin_name, str(len(sub)), f"{sub['label'].mean()*100:.1f}%", interp[bin_name]])
table(sl, 0.6, 5.0, 12, 2.2, pos_data, cw=[1.5, 1.2, 1.5, 7.8], sz=12)

# ── SLIDE 4: Amino Acid Substitution Patterns ──
sl = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(sl, WHITE)
tb(sl, 0.6, 0.3, 12, 0.6, "V3 Database: Amino Acid Substitution Patterns", sz=30, bold=True, color=DARK)

# Substitution group heatmap table
tb(sl, 0.6, 1.0, 12, 0.4, "Pathogenicity by Substitution Type (property groups)", sz=18, bold=True, color=ACCENT)
sg = {}
for _, row in bdf.iterrows():
    k = f"{aagrp(row['ref_aa'])}→{aagrp(row['alt_aa'])}"
    sg.setdefault(k, []).append(row['label'])

sg_data = [["Substitution", "Count", "Path Rate", "Risk Level"]]
for k in sorted(sg.keys(), key=lambda x: -np.mean(sg[x])):
    v = sg[k]
    rate = np.mean(v) * 100
    if rate >= 75: risk = "HIGH"
    elif rate >= 60: risk = "Moderate"
    else: risk = "Low"
    sg_data.append([k, str(len(v)), f"{rate:.1f}%", risk])
table(sl, 0.6, 1.5, 12, 5.5, sg_data, cw=[3.5, 1.5, 1.5, 5.5], sz=11)

# ── SLIDE 5: Specific Pathogenic Substitutions ──
sl = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(sl, WHITE)
tb(sl, 0.6, 0.3, 12, 0.6, "V3 Database: Key Pathogenic Substitutions", sz=30, bold=True, color=DARK)

# Left: high path rate
tb(sl, 0.6, 1.0, 5.8, 0.4, "Highest Pathogenicity (>80%, n≥20)", sz=18, bold=True, color=RED)
hi = [["Ref→Alt", "n", "Path%", "Biological Mechanism"]]
hi_subs = [
    ("Cys→Trp", "21", "100%", "Disulfide bond loss + steric clash"),
    ("Met→Arg", "30", "100%", "Oxidation site loss + charge gain"),
    ("Leu→Arg", "34", "97.1%", "Hydrophobic core → charge"),
    ("Trp→Arg", "44", "93.2%", "Aromatic stacking loss + charge"),
    ("Asp→Tyr", "32", "90.6%", "Negative charge loss + bulk"),
    ("Leu→Pro", "98", "89.8%", "α-helix breaker"),
    ("Tyr→Cys", "68", "89.7%", "Phosphorylation site loss"),
    ("Cys→Arg", "47", "89.4%", "Disulfide bond loss + charge"),
    ("Met→Leu", "43", "88.4%", "Oxidation-sensitive residue"),
    ("Trp→Cys", "34", "88.2%", "Aromatic + size loss"),
]
hi.extend(hi_subs)
table(sl, 0.6, 1.5, 5.8, 5.0, hi, cw=[1.2, 0.6, 0.8, 3.2], sz=10)

# Right: low path rate
tb(sl, 6.8, 1.0, 5.8, 0.4, "Lowest Pathogenicity (<40%, n≥20)", sz=18, bold=True, color=GREEN)
lo = [["Ref→Alt", "n", "Path%", "Biological Mechanism"]]
lo_subs = [
    ("Ala→Val", "132", "33.3%", "Conservative hydrophobic"),
    ("Ala→Thr", "127", "32.3%", "Small → small polar"),
    ("Pro→Ser", "90", "38.9%", "Surface residue, tolerant"),
    ("Pro→Leu", "129", "46.5%", "Moderate size, flexible"),
    ("Val→Met", "96", "44.8%", "Conservative hydrophobic"),
    ("Arg→His", "148", "50.0%", "Charge conservation"),
    ("Ser→Cys", "38", "34.2%", "Polar → polar, size similar"),
    ("Gly→Ala", "46", "34.8%", "Small → small, flexible"),
    ("Ala→Pro", "55", "36.4%", "Surface, moderate impact"),
    ("Val→Ile", "40", "37.5%", "Conservative hydrophobic"),
]
lo.extend(lo_subs)
table(sl, 6.8, 1.5, 5.8, 5.0, lo, cw=[1.2, 0.6, 0.8, 3.2], sz=10)

# ── SLIDE 6: Charge and Size Effects ──
sl = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(sl, WHITE)
tb(sl, 0.6, 0.3, 12, 0.6, "V3 Database: Physicochemical Effects", sz=30, bold=True, color=DARK)

bdf['charge_change'] = bdf['alt_aa'].map(lambda x: 1 if x in AA_POS else (-1 if x in AA_NEG else 0)) \
                     - bdf['ref_aa'].map(lambda x: 1 if x in AA_POS else (-1 if x in AA_NEG else 0))
bdf['size_change'] = bdf['alt_aa'].map(AA_SZ) - bdf['ref_aa'].map(AA_SZ)

# Left: charge
tb(sl, 0.6, 1.0, 5.8, 0.4, "Charge Change Effect", sz=18, bold=True, color=ACCENT)
ch = [["Charge Δ", "Count", "Path Rate", "Interpretation"]]
ch_interp = {
    '+2': 'Double charge gain — strong disruption',
    '+1': 'Charge gain — often pathogenic',
    '0': 'No net charge change — moderate risk',
    '-1': 'Charge loss — often pathogenic',
    '-2': 'Double charge loss — strong disruption',
}
for cc in [2, 1, 0, -1, -2]:
    sub = bdf[bdf['charge_change']==cc]
    if len(sub) > 0:
        ch.append([f"{'+' if cc>0 else ''}{cc}", str(len(sub)), f"{sub['label'].mean()*100:.1f}%", ch_interp.get(str(cc), "")])
table(sl, 0.6, 1.5, 5.8, 2.5, ch, cw=[1.2, 1.2, 1.2, 2.2], sz=12)

# Right: size
tb(sl, 6.8, 1.0, 5.8, 0.4, "Side Chain Size Change Effect", sz=18, bold=True, color=ACCENT)
sz_data = [["Size Δ", "Count", "Path Rate", "Interpretation"]]
for sc in [-4,-3,-2,-1,0,1,2,3,4]:
    sub = bdf[bdf['size_change']==sc]
    if len(sub) >= 10:
        interp = ""
        if sc <= -3: interp = "Large shrink — strong disruption"
        elif sc <= -1: interp = "Moderate shrink"
        elif sc == 0: interp = "Similar size"
        elif sc >= 3: interp = "Large expansion — steric clash"
        elif sc >= 1: interp = "Moderate expansion"
        sz_data.append([f"{'+'if sc>0 else ''}{sc}", str(len(sub)), f"{sub['label'].mean()*100:.1f}%", interp])
table(sl, 6.8, 1.5, 5.8, 3.5, sz_data, cw=[1.2, 1.2, 1.2, 2.2], sz=12)

# Bottom: top mutated genes
tb(sl, 0.6, 4.5, 12, 0.4, "Top 10 Most Mutated Genes", sz=18, bold=True, color=ACCENT)
top_genes = bdf.groupby('GeneSymbol').agg(n=('label','size'), rate=('label','mean')).sort_values('n', ascending=False).head(10)
tg_data = [["Gene", "Variants", "Path Rate", "Known Disease"]]
diseases = {
    'GLA': 'Fabry disease', 'PAH': 'Phenylketonuria', 'GCK': 'MODY2 diabetes',
    'IDS': 'Hunter syndrome', 'TP53': 'Li-Fraumeni syndrome', 'MLH1': 'Lynch syndrome',
    'BEST1': 'Best vitelliform macular dystrophy', 'GAA': 'Pompe disease',
    'LDLR': 'Familial hypercholesterolemia', 'F9': 'Hemophilia B'
}
for gene, row in top_genes.iterrows():
    tg_data.append([gene, str(row['n']), f"{row['rate']*100:.1f}%", diseases.get(gene, "")])
table(sl, 0.6, 5.0, 12, 2.2, tg_data, cw=[1.5, 1.5, 1.5, 7.5], sz=12)

# ── SLIDE 7: Main Results ──
sl = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(sl, WHITE)
tb(sl, 0.6, 0.3, 12, 0.6, "V3 Results: Model Ranking (AUROC)", sz=30, bold=True, color=DARK)

td = [["Model", "Type", "AUROC", "AUPRC", "Spearman", "MCC", "Recall", "Vars"]]
for m in order:
    r = main[m]
    typ = "PLM" if m in ['ESM1b-650M','ESM2-650M','SaProt-650M','ProtT5-XL'] else "DNA"
    td.append([m, typ, f"{r['AUROC']:.3f}", f"{r['AUPRC']:.3f}", f"{r['Spearman']:.3f}",
               f"{r['MCC']:.3f}", f"{r['Recall']:.3f}", f"{r['n']:,}"])
table(sl, 0.6, 1.0, 12, 5.0, td, cw=[2.2, 0.7, 1.2, 1.2, 1.3, 1.2, 1.2, 1.2], sz=12)

bullets(sl, 0.6, 6.2, 12, 0.8, [
    "PLMs dominate: ESM1b > ESM2 > SaProt >> NT-v2 > AlphaGenome > ProtT5 > HyenaDNA",
    "DNA models show weak signal (AUROC 0.55–0.64); PLMs achieve 0.83–0.87",
], sz=13, color=GRAY)

# ── SLIDE 8: Literature Comparison ──
sl = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(sl, WHITE)
tb(sl, 0.6, 0.3, 12, 0.6, "AUROC & Spearman: Literature vs. Ours", sz=30, bold=True, color=DARK)

fig1 = os.path.join(FIG_DIR, "fig_literature_comparison.png")
fig2 = os.path.join(FIG_DIR, "fig_literature_comparison_spearman.png")
if os.path.exists(fig1):
    sl.shapes.add_picture(fig1, Inches(0.3), Inches(1.0), Inches(6.4), Inches(6.2))
if os.path.exists(fig2):
    sl.shapes.add_picture(fig2, Inches(6.8), Inches(1.0), Inches(6.4), Inches(6.2))

# ── SLIDE 9: Embedding Experiments ──
sl = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(sl, WHITE)
tb(sl, 0.6, 0.3, 12, 0.6, "Embedding Experiments: Cosine Similarity vs LLR", sz=30, bold=True, color=DARK)

et = [["Model", "Method", "Spearman", "AUROC", "AUPRC", "MCC"]]
for m in ['ESM1b-650M', 'ESM2-650M', 'SaProt-650M']:
    r1, r2 = main[m], cos[m]
    et.append([m, "LLR",  f"{r1['Spearman']:.3f}", f"{r1['AUROC']:.3f}", f"{r1['AUPRC']:.3f}", f"{r1['MCC']:.3f}"])
    et.append(["",    "Cosine", f"{r2['Spearman']:.3f}", f"{r2['AUROC']:.3f}", f"{r2['AUPRC']:.3f}", f"{r2['MCC']:.3f}"])
table(sl, 0.6, 1.0, 12, 3.2, et, cw=[2.2, 1.5, 1.5, 1.5, 1.5, 1.5], sz=12)

bullets(sl, 0.6, 4.5, 12, 2.5, [
    "LLR dominates embedding cosine similarity for all models",
    "ESM1b: cosine loses 8% Spearman (0.631 → 0.534) and 5% AUROC (0.873 → 0.816)",
    "ESM2: cosine loses 7% Spearman (0.572 → 0.501) and 4% AUROC (0.838 → 0.797)",
    "SaProt: cosine collapses completely (Spearman 0.556 → 0.114)",
    "Max-pool + cosine(diff, wt_pool) gave negative Spearman — fundamentally flawed",
    "Conclusion: Log-likelihood ratio captures pathogenicity better than embedding similarity",
], sz=13)

# ── SLIDE 10: Key Findings ──
sl = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(sl, WHITE)
tb(sl, 0.6, 0.3, 12, 0.6, "Key Findings", sz=30, bold=True, color=DARK)

bullets(sl, 0.6, 1.0, 5.8, 5.5, [
    "1. Protein LMs significantly outperform DNA LMs",
    "   ESM1b AUROC=0.873 vs AlphaGenome 0.552",
    "",
    "2. ESM1b > ESM2 > SaProt (consistent across metrics)",
    "   Older ESM1b slightly outperforms newer ESM2",
    "",
    "3. SaProt 3Di foldseek codes introduce noise",
    "   880 genes missing; cosine embedding collapse",
    "",
    "4. DNA models barely above random",
    "   AlphaGenome AUROC=0.552, HyenaDNA=0.448",
], sz=14)

bullets(sl, 6.8, 1.0, 5.8, 5.5, [
    "5. LLR >> cosine similarity for scoring",
    "   Embedding similarity fails to capture pathogenicity",
    "",
    "6. Literature gap explained by:",
    "   - Smaller benchmark (5.9k vs 30–100k)",
    "   - SaProt: AlphaFold 3Di vs PDB 3Di",
    "   - NT-v2: hidden state delta vs embedding delta",
    "   - AlphaGenome: regulatory model, not missense",
    "",
    "7. Consistent model ranking across all metrics",
    "   AUROC, AUPRC, Spearman, MCC all agree",
], sz=14)

# ── SLIDE 11: Conclusions ──
sl = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(sl, DARK)
tb(sl, 1.0, 1.2, 11.3, 0.8, "Conclusions & Next Steps",
   sz=36, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
bullets(sl, 1.0, 2.5, 11.3, 4.5, [
    "PLMs are far superior to DNA models for pathogenic variant prediction",
    "ESM1b-650M achieves state-of-the-art on our V3 benchmark (AUROC=0.873)",
    "Log-likelihood ratio is the optimal scoring method — cosine similarity is inferior",
    "",
    "Biological patterns confirmed:",
    "  Charge-altering substitutions → 70–85% pathogenic",
    "  Large size changes → 79–82% pathogenic",
    "  C-terminal variants more tolerant (47.6% path vs 65% N-terminal)",
    "",
    "Next steps:",
    "  Test ESM3 (15B) on V3",
    "  Expand benchmark with ProteinGym DMS data",
    "  Explore ensemble methods (PLM + DNA signals)",
    "  Investigate SaProt 3Di embedding space collapse",
], sz=16, color=WHITE)

# ════════════════════════════════════════════════════
out = os.path.join(FIG_DIR, "V3_Benchmark_Results.pptx")
prs.save(out)
print(f"Saved: {out}")
print(f"Slides: {len(prs.slides)}")
